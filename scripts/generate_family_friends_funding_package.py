from __future__ import annotations

import math
import re
import shutil
import textwrap
from dataclasses import dataclass
from pathlib import Path
from typing import Dict, Iterable, List, Optional, Tuple

import matplotlib

matplotlib.use("Agg")
import matplotlib.pyplot as plt
import numpy as np
import pandas as pd
from matplotlib.backends.backend_pdf import PdfPages
from matplotlib.ticker import FuncFormatter
from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

from live.replay_manifest import write_run_manifest


ROOT = Path(__file__).resolve().parents[1]
OUT = ROOT / "live" / "state" / "family_friends_funding_package"
CHARTS = OUT / "charts"
EXHIBITS = OUT / "exhibits"
FAIR_BENCHMARK_DIR = ROOT / "mnq" / "case_studies" / "fair_benchmark_comparison"

RUNWAY_MONTHS = 12
DEV_STIPEND_MONTHLY = 3_000.0
DATABENTO_LIVE_MONTHLY = 179.0
STARTER_ASK = 145_000.0
MNQ_MINIMUM = 50_000.0
NQ_MINIMUM = 250_000.0

PETIT_NAME = "Nasdaq Intraday Mini"
GRAND_NAME = "Nasdaq Intraday"
PROGRAM_NAME = "Nasdaq Futures Diversification Platform"
PITCH_PROGRAM_NAME = "Nasdaq Intraday Program"
FOUNDER_PARTNER_NAME = "Founder Partner Accounts"
FOUNDER_PARTNER_MGMT_FEE = 0.0
FOUNDER_PARTNER_PERFORMANCE_FEE = 10.0
MODEL_AUDIT_FEE_PER_CLOSED_UNIT = 1.50
STONEX_STANDARD_PER_SIDE = 1.29
STONEX_MICRO_PER_SIDE = 0.50
STONEX_SIM_SOFTWARE_MONTHLY = 49.95
NFA_CTA_FORM_7R_FEE = 200.0
NFA_CTA_DUES = 750.0
NFA_AP_OR_PRINCIPAL_FEE = 85.0
SERIES_3_EXAM_FEE = 140.0
NFA_SERIES3_REGISTRATION_ADMIN_ALLOWANCE = 2_500.0
CTA_LEGAL_COMPLIANCE_AUDIT_RESERVE = 30_000.0
BUDGET_EXHIBIT = "runway_budget.csv"
PDF_DIR = OUT / "pdfs"
POWERPOINT_PATH = OUT / "PITCH_DECK.pptx"

PUBLIC_DOCS = [
    "PITCH_DECK.md",
    "SHORT_MEMO.md",
    "COST_AND_RUNWAY.md",
    "REPORTING_TEMPLATES.md",
    "RISK_AND_DISCLOSURE_NOTES.md",
]


@dataclass(frozen=True)
class StrategySource:
    slug: str
    public_name: str
    tier: str
    market: str
    minimum_capital: float
    yearly_path: Path
    equity_path: Path
    unit_trades_path: Path
    qqq_path: Path
    summary_path: Path
    preferred: bool = False
    use_for_deck: bool = False


SOURCES = [
    StrategySource(
        slug="nq_prior_opposed_full",
        public_name="Prior-Opposed Gated Intraday System",
        tier=GRAND_NAME,
        market="NQ",
        minimum_capital=NQ_MINIMUM,
        yearly_path=ROOT / "live/state/nq_v2b_prior_opposed_stpmc_full_history_raw/yearly_breakdown.csv",
        equity_path=ROOT
        / "live/state/nq_v2b_prior_opposed_stpmc_full_history_raw/states/nq_v2b_prior_opposed_stpmc_only_S_1_1_3/equity_curve.csv",
        unit_trades_path=ROOT
        / "live/state/nq_v2b_prior_opposed_stpmc_full_history_raw/states/nq_v2b_prior_opposed_stpmc_only_S_1_1_3/unit_trades.csv",
        qqq_path=ROOT / "data/benchmarks/QQQ_2010-06-06_2026-03-08_yahoo_daily.csv",
        summary_path=ROOT / "live/state/nq_v2b_prior_opposed_stpmc_full_history_raw/summary.csv",
        preferred=True,
        use_for_deck=True,
    ),
    StrategySource(
        slug="mnq_prior_opposed",
        public_name="Prior-Opposed Gated Intraday System",
        tier=PETIT_NAME,
        market="MNQ",
        minimum_capital=MNQ_MINIMUM,
        yearly_path=ROOT / "live/state/mnq_v2b_prior_opposed_stpmc_broker_like/robustness_audit/yearly_breakdown.csv",
        equity_path=ROOT
        / "live/state/mnq_v2b_prior_opposed_stpmc_broker_like/states/mnq_v2b_prior_opposed_stpmc_only_S_1_1_3/equity_curve.csv",
        unit_trades_path=ROOT
        / "live/state/mnq_v2b_prior_opposed_stpmc_broker_like/states/mnq_v2b_prior_opposed_stpmc_only_S_1_1_3/unit_trades.csv",
        qqq_path=ROOT / "data/benchmarks/QQQ_2021-03-04_2026-03-06_yahoo_daily.csv",
        summary_path=ROOT / "live/state/mnq_v2b_prior_opposed_stpmc_broker_like/summary.csv",
        preferred=True,
        use_for_deck=True,
    ),
    StrategySource(
        slug="nq_ungated_intraday",
        public_name="Ungated Intraday Breakout System",
        tier=GRAND_NAME,
        market="NQ",
        minimum_capital=NQ_MINIMUM,
        yearly_path=ROOT / "live/state/v2b_sizing_sweep/nq_1_1_3_yearly.csv",
        equity_path=ROOT / "live/state/v2b_sizing_sweep/states/nq_v2b_sizing_S_1_1_3/equity_curve.csv",
        unit_trades_path=ROOT / "live/state/v2b_sizing_sweep/states/nq_v2b_sizing_S_1_1_3/unit_trades.csv",
        qqq_path=ROOT / "data/benchmarks/QQQ_2021-03-04_2026-03-06_yahoo_daily.csv",
        summary_path=ROOT / "live/state/v2b_sizing_sweep/summary_partial.csv",
        use_for_deck=True,
    ),
    StrategySource(
        slug="mnq_ungated_intraday",
        public_name="Ungated Intraday Breakout System",
        tier=PETIT_NAME,
        market="MNQ",
        minimum_capital=MNQ_MINIMUM,
        yearly_path=ROOT / "live/state/v2b_sizing_sweep/mnq_1_1_3_yearly.csv",
        equity_path=ROOT / "live/state/v2b_sizing_sweep/states/mnq_v2b_sizing_S_1_1_3/equity_curve.csv",
        unit_trades_path=ROOT / "live/state/v2b_sizing_sweep/states/mnq_v2b_sizing_S_1_1_3/unit_trades.csv",
        qqq_path=ROOT / "data/benchmarks/QQQ_2021-03-04_2026-03-06_yahoo_daily.csv",
        summary_path=ROOT / "live/state/v2b_sizing_sweep/summary_partial.csv",
        use_for_deck=True,
    ),
]


def money(value: float, digits: int = 0) -> str:
    return f"${value:,.{digits}f}"


def pct(value: float, digits: int = 1) -> str:
    return f"{value:.{digits}f}%"


def ratio(value: float, digits: int = 2) -> str:
    if value is None or not np.isfinite(value):
        return "n/a"
    return f"{value:.{digits}f}"


def md_table(rows: List[Dict[str, object]], columns: List[str]) -> str:
    if not rows:
        return ""
    lines = [
        "| " + " | ".join(columns) + " |",
        "| " + " | ".join(["---"] * len(columns)) + " |",
    ]
    for row in rows:
        lines.append("| " + " | ".join(str(row.get(col, "")) for col in columns) + " |")
    return "\n".join(lines)


def read_yearly(src: StrategySource) -> pd.DataFrame:
    df = pd.read_csv(src.yearly_path)
    year_col = "year" if "year" in df.columns else "Year"
    net_col = "net_usd" if "net_usd" in df.columns else "Net"
    trades_col = "trades" if "trades" in df.columns else ("campaigns" if "campaigns" in df.columns else "Trades")
    win_col = "win_rate_pct" if "win_rate_pct" in df.columns else "Win %"
    pf_col = "profit_factor" if "profit_factor" in df.columns else "PF"
    dd_col = (
        "campaign_closed_dd_usd"
        if "campaign_closed_dd_usd" in df.columns
        else ("closed_dd_usd" if "closed_dd_usd" in df.columns else "MTM Stress DD")
    )
    out = pd.DataFrame(
        {
            "year": df[year_col].astype(int),
            "trades": df[trades_col].astype(float),
            "net_usd": df[net_col].astype(float),
            "win_rate_pct": df[win_col].astype(float),
            "profit_factor": df[pf_col].astype(float),
            "drawdown_usd": df[dd_col].astype(float),
        }
    )
    out["strategy_slug"] = src.slug
    out["public_name"] = src.public_name
    out["tier"] = src.tier
    out["market"] = src.market
    out["return_on_minimum_pct"] = out["net_usd"] / src.minimum_capital * 100.0
    return out.sort_values("year")


def load_strategy_daily(src: StrategySource) -> pd.DataFrame:
    df = pd.read_csv(src.equity_path, usecols=["ts", "close_equity_usd", "intrabar_stress_equity_usd"])
    ts = pd.to_datetime(df["ts"], utc=True, errors="coerce").dt.tz_convert("America/New_York")
    df = df[ts.notna()].copy()
    df["date"] = ts[ts.notna()].dt.date
    daily = (
        df.groupby("date", as_index=False)
        .agg(
            close_equity_usd=("close_equity_usd", "last"),
            intrabar_stress_equity_usd=("intrabar_stress_equity_usd", "min"),
        )
        .copy()
    )
    daily["date"] = pd.to_datetime(daily["date"])
    daily["strategy_equity"] = src.minimum_capital + daily["close_equity_usd"].astype(float)
    daily["stress_equity"] = src.minimum_capital + daily["intrabar_stress_equity_usd"].astype(float)
    daily["strategy_return"] = daily["strategy_equity"].pct_change().fillna(0.0)
    daily["strategy_drawdown_pct"] = daily["strategy_equity"] / daily["strategy_equity"].cummax() - 1.0
    daily["stress_drawdown_pct"] = daily["stress_equity"] / daily["strategy_equity"].cummax() - 1.0
    return daily[
        [
            "date",
            "strategy_equity",
            "stress_equity",
            "strategy_return",
            "strategy_drawdown_pct",
            "stress_drawdown_pct",
            "close_equity_usd",
            "intrabar_stress_equity_usd",
        ]
    ].sort_values("date")


def load_qqq_daily(
    src: StrategySource,
    start_date: Optional[pd.Timestamp] = None,
    end_date: Optional[pd.Timestamp] = None,
) -> pd.DataFrame:
    qqq = pd.read_csv(src.qqq_path)
    qqq["date"] = pd.to_datetime(qqq["date"] if "date" in qqq.columns else qqq["Date"])
    price_col = "adj_close" if "adj_close" in qqq.columns else ("Adj Close" if "Adj Close" in qqq.columns else "close")
    qqq = qqq.sort_values("date")
    if start_date is not None:
        qqq = qqq[qqq["date"] >= pd.to_datetime(start_date)]
    if end_date is not None:
        qqq = qqq[qqq["date"] <= pd.to_datetime(end_date)]
    if qqq.empty:
        raise ValueError(f"No QQQ benchmark rows for {src.slug} after applying strategy-window clipping")
    qqq["qqq_equity"] = src.minimum_capital * qqq[price_col].astype(float) / float(qqq[price_col].iloc[0])
    qqq["qqq_return"] = qqq["qqq_equity"].pct_change().fillna(0.0)
    qqq["qqq_drawdown_pct"] = qqq["qqq_equity"] / qqq["qqq_equity"].cummax() - 1.0
    return qqq[["date", "qqq_equity", "qqq_return", "qqq_drawdown_pct"]]


def qqq_annual_returns(src: StrategySource) -> pd.DataFrame:
    strategy_daily = load_strategy_daily(src)
    qqq = load_qqq_daily(src, strategy_daily["date"].min(), strategy_daily["date"].max())
    qqq["year"] = qqq["date"].dt.year
    first = qqq.groupby("year").first()["qqq_equity"]
    last = qqq.groupby("year").last()["qqq_equity"]
    out = ((last / first - 1.0) * 100.0).rename("qqq_return_pct").reset_index()
    return out


def official_summary(src: StrategySource) -> Dict[str, float]:
    if not src.summary_path.exists():
        return {}
    summary = pd.read_csv(src.summary_path)
    if "strategy_id" in summary.columns and len(summary) == 1:
        row = summary.iloc[0]
    elif "slug" in summary.columns and "market" in summary.columns:
        rows = summary[
            (summary["market"].astype(str).str.lower() == src.market.lower())
            & (summary["slug"].astype(str) == "S_1_1_3")
        ]
        if rows.empty:
            return {}
        row = rows.iloc[0]
    else:
        return {}

    return {
        "net_usd": float(row["net_usd"]) if "net_usd" in row else float("nan"),
        "closed_dd_usd": float(row["closed_dd_usd"]) if "closed_dd_usd" in row else float("nan"),
        "intrabar_stress_dd_usd": float(row["intrabar_stress_dd_usd"])
        if "intrabar_stress_dd_usd" in row
        else float("nan"),
        "win_rate_pct": float(row["win_rate_pct"]) if "win_rate_pct" in row else float("nan"),
        "profit_factor": float(row["profit_factor"]) if "profit_factor" in row else float("nan"),
        "net_over_stress": float(row["net_over_stress"] if "net_over_stress" in row else row.get("net_over_stress_dd", np.nan)),
    }


def stonex_fee_estimate(src: StrategySource) -> Dict[str, float]:
    unit_trades = pd.read_csv(src.unit_trades_path, usecols=["trade_id"])
    closed_units = int(len(unit_trades))
    per_side = STONEX_MICRO_PER_SIDE if src.market.upper().startswith("M") else STONEX_STANDARD_PER_SIDE
    round_turn = per_side * 2.0
    stonex_commission = closed_units * round_turn
    model_audit_fee = closed_units * MODEL_AUDIT_FEE_PER_CLOSED_UNIT
    return {
        "closed_units": closed_units,
        "stonex_per_side": per_side,
        "stonex_round_turn": round_turn,
        "stonex_commission_usd": stonex_commission,
        "model_audit_fee_usd": model_audit_fee,
        "stonex_delta_vs_model_fee_usd": stonex_commission - model_audit_fee,
    }


def conservative_average(yearly: pd.DataFrame) -> Tuple[float, float, pd.DataFrame, pd.DataFrame]:
    top = yearly.nlargest(min(3, len(yearly)), "net_usd")
    remainder = yearly.drop(top.index)
    avg_net = float(remainder["net_usd"].mean()) if not remainder.empty else float("nan")
    avg_return = float(remainder["return_on_minimum_pct"].mean()) if not remainder.empty else float("nan")
    return avg_net, avg_return, top.sort_values("year"), remainder.sort_values("year")


def summarize_strategy(src: StrategySource) -> Dict[str, object]:
    yearly = read_yearly(src)
    official = official_summary(src)
    fees = stonex_fee_estimate(src)
    qqq = qqq_annual_returns(src)
    strategy_daily = load_strategy_daily(src)
    qqq_daily = load_qqq_daily(src, strategy_daily["date"].min(), strategy_daily["date"].max())
    daily = strategy_daily.merge(qqq_daily, on="date", how="inner")
    avg_net, avg_return, top, remainder = conservative_average(yearly)
    qqq_join = yearly.merge(qqq, on="year", how="left")
    qqq_down = qqq_join[qqq_join["qqq_return_pct"] < 0]
    qqq_up = qqq_join[qqq_join["qqq_return_pct"] >= 0]
    start = daily["date"].iloc[0]
    end = daily["date"].iloc[-1]
    derived_stress_dd = float((daily["stress_equity"] - daily["strategy_equity"].cummax()).min())
    stress_dd = official.get("intrabar_stress_dd_usd", derived_stress_dd)
    return {
        "slug": src.slug,
        "public_name": src.public_name,
        "tier": src.tier,
        "market": src.market,
        "minimum_capital": src.minimum_capital,
        "start_date": start,
        "end_date": end,
        "years": len(yearly),
        "total_net": float(yearly["net_usd"].sum()),
        "total_return_on_minimum_pct": float(yearly["net_usd"].sum() / src.minimum_capital * 100.0),
        "avg_all_years_net": float(yearly["net_usd"].mean()),
        "avg_all_years_return_pct": float(yearly["return_on_minimum_pct"].mean()),
        "advertised_avg_net": avg_net,
        "advertised_avg_return_pct": avg_return,
        "removed_top3_years": ", ".join(str(int(y)) for y in top["year"]),
        "removed_top3_net": float(top["net_usd"].sum()),
        "daily_corr_to_qqq": float(daily["strategy_return"].corr(daily["qqq_return"])),
        "avg_return_in_qqq_down_years_pct": float(qqq_down["return_on_minimum_pct"].mean())
        if not qqq_down.empty
        else float("nan"),
        "avg_return_in_qqq_up_years_pct": float(qqq_up["return_on_minimum_pct"].mean())
        if not qqq_up.empty
        else float("nan"),
        "qqq_down_years": ", ".join(str(int(y)) for y in qqq_down["year"]),
        "max_stress_dd_usd": stress_dd,
        "max_stress_dd_pct_of_minimum": float(stress_dd / src.minimum_capital * 100.0),
        "peak_to_trough_stress_pct": float(daily["stress_drawdown_pct"].min() * 100.0),
        **fees,
    }


def build_annual_comparison(summaries: Dict[str, Dict[str, object]]) -> pd.DataFrame:
    rows = []
    for src in SOURCES:
        yearly = read_yearly(src)
        qqq = qqq_annual_returns(src)
        joined = yearly.merge(qqq, on="year", how="left")
        for row in joined.itertuples(index=False):
            rows.append(
                {
                    "strategy_slug": src.slug,
                    "strategy": src.public_name,
                    "tier": src.tier,
                    "market": src.market,
                    "year": int(row.year),
                    "strategy_net_usd": float(row.net_usd),
                    "strategy_return_on_minimum_pct": float(row.return_on_minimum_pct),
                    "qqq_return_pct": float(row.qqq_return_pct) if np.isfinite(row.qqq_return_pct) else np.nan,
                    "qqq_year_type": "QQQ down" if np.isfinite(row.qqq_return_pct) and row.qqq_return_pct < 0 else "QQQ up",
                    "win_rate_pct": float(row.win_rate_pct),
                    "profit_factor": float(row.profit_factor),
                    "drawdown_usd": float(row.drawdown_usd),
                }
            )
    return pd.DataFrame(rows)


def ensure_dirs() -> None:
    if OUT.exists():
        shutil.rmtree(OUT)
    CHARTS.mkdir(parents=True, exist_ok=True)
    EXHIBITS.mkdir(parents=True, exist_ok=True)
    PDF_DIR.mkdir(parents=True, exist_ok=True)


def write_exhibits(summaries: Dict[str, Dict[str, object]]) -> None:
    pd.DataFrame(summaries.values()).to_csv(EXHIBITS / "strategy_summary.csv", index=False)
    build_annual_comparison(summaries).to_csv(EXHIBITS / "annual_returns_vs_qqq.csv", index=False)
    stonex_rows = []
    for src in SOURCES:
        s = summaries[src.slug]
        stonex_rows.append(
            {
                "strategy_slug": src.slug,
                "strategy": src.public_name,
                "tier": src.tier,
                "market": src.market,
                "closed_units": int(s["closed_units"]),
                "stonex_per_side": float(s["stonex_per_side"]),
                "stonex_round_turn": float(s["stonex_round_turn"]),
                "stonex_commission_usd": float(s["stonex_commission_usd"]),
                "model_audit_fee_usd": float(s["model_audit_fee_usd"]),
                "stonex_delta_vs_model_fee_usd": float(s["stonex_delta_vs_model_fee_usd"]),
            }
        )
    pd.DataFrame(stonex_rows).to_csv(EXHIBITS / "stonex_fee_estimate.csv", index=False)

    budget = pd.DataFrame(
        [
            {
                "Category": "Founder development stipend",
                "Basis": f"$3,000/month x {RUNWAY_MONTHS} months",
                "Amount": DEV_STIPEND_MONTHLY * RUNWAY_MONTHS,
            },
            {"Category": "Segregated MNQ research/test capital", "Basis": "Held as strategy test capital", "Amount": 50_000.0},
            {
                "Category": "Databento live data estimate",
                "Basis": f"$179/month x {RUNWAY_MONTHS} months",
                "Amount": DATABENTO_LIVE_MONTHLY * RUNWAY_MONTHS,
            },
            {
                "Category": "StoneX simulation software",
                "Basis": f"$49.95/month x {RUNWAY_MONTHS} months, per provided fee schedule",
                "Amount": round(STONEX_SIM_SOFTWARE_MONTHLY * RUNWAY_MONTHS, 2),
            },
            {"Category": "Cloud/runtime/monitoring", "Basis": "EC2, logs, alerts, backups", "Amount": 2_000.0},
            {
                "Category": "Broker/API/exchange/commission/slippage allowance",
                "Basis": "StoneX commissions, exchange/regulatory fees, wires, and slippage buffer",
                "Amount": 10_000.0,
            },
            {
                "Category": "NFA/Series 3/registration admin",
                "Basis": "Series 3, Form 7-R, AP/principal, CTA dues, fingerprint/admin buffer",
                "Amount": NFA_SERIES3_REGISTRATION_ADMIN_ALLOWANCE,
            },
            {
                "Category": "CTA disclosure/legal/compliance/audit reserve",
                "Basis": "CTA registration-first counsel, disclosure document, NFA review responses, audit/accounting setup",
                "Amount": CTA_LEGAL_COMPLIANCE_AUDIT_RESERVE,
            },
            {"Category": "Contingency", "Basis": "Unexpected data, infrastructure, and compliance costs", "Amount": 10_000.0},
        ]
    )
    budget.to_csv(EXHIBITS / BUDGET_EXHIBIT, index=False)

    regulatory_costs = pd.DataFrame(
        [
            {
                "Item": "NFA CTA Form 7-R application fee",
                "Assumption": money(NFA_CTA_FORM_7R_FEE, 0),
                "Source": "NFA CTA registration requirements",
                "URL": "https://www.nfa.futures.org/registration-membership/who-has-to-register/cta.html",
            },
            {
                "Item": "NFA CTA initial/annual dues",
                "Assumption": money(NFA_CTA_DUES, 0),
                "Source": "NFA Membership Dues and Fees",
                "URL": "https://www.nfa.futures.org/registration-membership/dues-revenue-structure.html",
            },
            {
                "Item": "Principal/AP application fee",
                "Assumption": money(NFA_AP_OR_PRINCIPAL_FEE, 0),
                "Source": "NFA CTA registration requirements",
                "URL": "https://www.nfa.futures.org/registration-membership/who-has-to-register/cta.html",
            },
            {
                "Item": "Series 3 exam fee",
                "Assumption": money(SERIES_3_EXAM_FEE, 0),
                "Source": "NFA proficiency/exam guidance",
                "URL": "https://www.nfa.futures.org/registration-membership/study-outlines/index.html",
            },
            {
                "Item": "CTA disclosure/legal/compliance/audit reserve",
                "Assumption": money(CTA_LEGAL_COMPLIANCE_AUDIT_RESERVE, 0),
                "Source": "Planning reserve based on attached CTA/CTC cost note; counsel quote required",
                "URL": "local attachment: Comprehensive CTA / CTC Launch & Operational Cost Summary",
            },
        ]
    )
    regulatory_costs.to_csv(EXHIBITS / "regulatory_cost_assumptions.csv", index=False)

    manifest = pd.DataFrame(
        [
            {
                "strategy_slug": src.slug,
                "public_name": src.public_name,
                "tier": src.tier,
                "yearly_source": str(src.yearly_path.relative_to(ROOT)),
                "equity_source": str(src.equity_path.relative_to(ROOT)),
                "unit_trades_source": str(src.unit_trades_path.relative_to(ROOT)),
                "qqq_source": str(src.qqq_path.relative_to(ROOT)),
                "summary_source": str(src.summary_path.relative_to(ROOT)),
            }
            for src in SOURCES
        ]
    )
    manifest.to_csv(EXHIBITS / "source_manifest_internal.csv", index=False)


def write_charts(summaries: Dict[str, Dict[str, object]]) -> None:
    annual = build_annual_comparison(summaries)

    fig, ax = plt.subplots(figsize=(11.5, 5.5))
    preferred = [s for s in SOURCES if s.preferred]
    colors = ["#1f77b4", "#2ca02c"]
    for idx, src in enumerate(preferred):
        strategy_daily = load_strategy_daily(src)
        d = strategy_daily.merge(load_qqq_daily(src, strategy_daily["date"].min(), strategy_daily["date"].max()), on="date", how="inner")
        d["strategy_index"] = d["strategy_equity"] / src.minimum_capital
        d["qqq_index"] = d["qqq_equity"] / src.minimum_capital
        label = f"{src.tier} strategy"
        color = colors[idx % len(colors)]
        ax.plot(d["date"], d["strategy_index"], linewidth=1.8, color=color, label=label)
        ax.plot(
            d["date"],
            d["qqq_index"],
            color=color,
            linestyle="--",
            linewidth=1.1,
            alpha=0.7,
            label=f"QQQ same window ({src.market})",
        )
    ax.set_title("Preferred Tier Equity Growth vs QQQ")
    ax.set_ylabel("Growth of $1.00")
    ax.grid(True, alpha=0.25)
    ax.legend(loc="upper left")
    fig.tight_layout()
    fig.savefig(CHARTS / "preferred_tiers_vs_qqq_equity.png", dpi=160)
    plt.close(fig)

    nq_src = next(s for s in SOURCES if s.slug == "nq_prior_opposed_full")
    nq_strategy_daily = load_strategy_daily(nq_src)
    nq = nq_strategy_daily.merge(
        load_qqq_daily(nq_src, nq_strategy_daily["date"].min(), nq_strategy_daily["date"].max()),
        on="date",
        how="inner",
    )
    nq["strategy_index"] = nq["strategy_equity"] / nq_src.minimum_capital
    nq["qqq_index"] = nq["qqq_equity"] / nq_src.minimum_capital
    fig, ax = plt.subplots(figsize=(11.5, 5.5))
    ax.plot(nq["date"], nq["strategy_index"], color="#1f77b4", linewidth=1.9, label="NQ prior-opposed strategy")
    ax.plot(nq["date"], nq["qqq_index"], color="#333333", linewidth=1.5, label="QQQ buy-and-hold")
    ax.set_title("NQ Prior-Opposed Full-History Equity Growth vs QQQ")
    ax.set_ylabel("Growth of $1.00")
    ax.grid(True, alpha=0.25)
    ax.legend(loc="upper left")
    fig.tight_layout()
    fig.savefig(CHARTS / "nq_prior_opposed_full_vs_qqq_equity.png", dpi=160)
    plt.close(fig)

    fig, ax = plt.subplots(figsize=(11.5, 5.5))
    deck_rows = annual[annual["strategy_slug"].isin([s.slug for s in SOURCES if s.use_for_deck])]
    pivot = deck_rows.pivot_table(
        index=["strategy", "tier", "market"],
        values="strategy_return_on_minimum_pct",
        aggfunc="mean",
    ).reset_index()
    labels = [f"{r.market} {r.strategy}\n{r.tier}" for r in pivot.itertuples(index=False)]
    ax.barh(labels, pivot["strategy_return_on_minimum_pct"], color="#386cb0")
    ax.set_title("Average Annual Return on Tier Minimum Capital")
    ax.set_xlabel("Average annual return (%)")
    ax.grid(axis="x", alpha=0.25)
    fig.tight_layout()
    fig.savefig(CHARTS / "average_annual_return_by_strategy.png", dpi=160)
    plt.close(fig)

    fig, ax = plt.subplots(figsize=(12, 5.5))
    selected = annual[annual["strategy_slug"].isin(["mnq_prior_opposed", "mnq_ungated_intraday"])].copy()
    years = sorted(selected["year"].unique())
    x = np.arange(len(years))
    width = 0.25
    for i, slug in enumerate(["mnq_prior_opposed", "mnq_ungated_intraday"]):
        vals = (
            selected[selected["strategy_slug"] == slug]
            .set_index("year")
            .reindex(years)["strategy_return_on_minimum_pct"]
            .fillna(0.0)
        )
        label = summaries[slug]["public_name"]
        ax.bar(x + (i - 0.5) * width, vals, width=width, label=label)
    qvals = selected.drop_duplicates("year").set_index("year").reindex(years)["qqq_return_pct"].fillna(0.0)
    ax.plot(x + width, qvals, color="#333333", marker="o", linewidth=1.5, label="QQQ")
    ax.axhline(0, color="#222222", linewidth=0.8)
    ax.set_xticks(x)
    ax.set_xticklabels([str(y) for y in years])
    ax.set_title("Nasdaq Intraday Mini Annual Returns vs QQQ")
    ax.set_ylabel("Return (%)")
    ax.grid(axis="y", alpha=0.25)
    ax.legend(loc="upper left")
    fig.tight_layout()
    fig.savefig(CHARTS / "mini_annual_returns_vs_qqq.png", dpi=160)
    plt.close(fig)

    fig, ax = plt.subplots(figsize=(9, 4.8))
    rows = pd.DataFrame(summaries.values())
    rows = rows[rows["slug"].isin([s.slug for s in SOURCES if s.use_for_deck])]
    labels = [f"{r.market} {r.public_name}" for r in rows.itertuples(index=False)]
    ax.barh(labels, rows["daily_corr_to_qqq"], color=np.where(rows["daily_corr_to_qqq"] < 0, "#2ca25f", "#fdae6b"))
    ax.axvline(0, color="#333333", linewidth=0.9)
    ax.set_title("Daily Return Correlation To QQQ")
    ax.set_xlabel("Correlation")
    ax.grid(axis="x", alpha=0.25)
    fig.tight_layout()
    fig.savefig(CHARTS / "daily_correlation_to_qqq.png", dpi=160)
    plt.close(fig)

    scaling_daily, scaling_rows, _ = scaling_10y_nq_vs_qqq()
    scaling_daily.to_csv(EXHIBITS / "scaling_10y_nq_prior_vs_qqq_50k_daily.csv", index=False)
    pd.DataFrame(scaling_rows).to_csv(EXHIBITS / "scaling_10y_nq_prior_vs_qqq_50k_summary.csv", index=False)
    fig, ax = plt.subplots(figsize=(10.4, 4.9))
    ax.plot(
        scaling_daily["date"],
        scaling_daily["prior_equity"],
        color="#006ba4",
        linewidth=2.1,
        label="NQ prior-opposed fixed 1 base book",
    )
    ax.plot(
        scaling_daily["date"],
        scaling_daily["qqq_equity"],
        color="#b35c1e",
        linewidth=1.8,
        label="QQQ fully invested",
    )
    ax.set_title("Fair Benchmark Scaling 10Y: $50k Equity Growth")
    ax.set_ylabel("Account equity")
    ax.yaxis.set_major_formatter(FuncFormatter(lambda x, _: f"${x/1000:,.0f}k"))
    ax.grid(True, alpha=0.22)
    ax.legend(loc="upper left", frameon=False)
    fig.tight_layout()
    fig.savefig(CHARTS / "scaling_10y_nq_prior_vs_qqq_50k.png", dpi=170)
    plt.close(fig)


def summary_rows(summaries: Dict[str, Dict[str, object]], preferred_only: bool = False) -> List[Dict[str, object]]:
    rows = []
    allowed = {s.slug for s in SOURCES if s.preferred} if preferred_only else {s.slug for s in SOURCES if s.use_for_deck}
    for src in SOURCES:
        if src.slug not in allowed:
            continue
        s = summaries[src.slug]
        rows.append(
            {
                "Tier": s["tier"],
                "System": s["public_name"],
                "Market": s["market"],
                "Minimum": money(float(s["minimum_capital"]), 0),
                "Window": f"{s['start_date'].date()} to {s['end_date'].date()}",
                "Total Net": money(float(s["total_net"]), 0),
                "Advertised Avg/Yr": money(float(s["advertised_avg_net"]), 0),
                "Advertised Return/Yr": pct(float(s["advertised_avg_return_pct"]), 1),
                "Corr to QQQ": ratio(float(s["daily_corr_to_qqq"]), 2),
            }
        )
    return rows


def annual_rows_for(src_slug: str) -> List[Dict[str, object]]:
    src = next(s for s in SOURCES if s.slug == src_slug)
    yearly = read_yearly(src)
    qqq = qqq_annual_returns(src)
    df = yearly.merge(qqq, on="year", how="left").sort_values("year")
    rows = []
    for row in df.itertuples(index=False):
        rows.append(
            {
                "Year": int(row.year),
                "System Return": pct(float(row.return_on_minimum_pct), 1),
                "Net": money(float(row.net_usd), 0),
                "QQQ Return": pct(float(row.qqq_return_pct), 1),
                "QQQ Regime": "QQQ down" if np.isfinite(row.qqq_return_pct) and row.qqq_return_pct < 0 else "QQQ up",
            }
        )
    return rows


def scaling_10y_nq_vs_qqq(initial_capital: float = 50_000.0) -> Tuple[pd.DataFrame, List[Dict[str, object]], pd.Series]:
    """Fair Benchmark Scaling 10Y exhibit for the Athena-style slide.

    This is intentionally the fixed one-base-book NQ study requested for the
    performance page. It is not the 3x-stress sizing recommendation.
    """

    prior_path = FAIR_BENCHMARK_DIR / "prior_opposed_intraday_nq_fixed_1book_50k_daily.csv"
    qqq_path = ROOT / "data" / "benchmarks" / "QQQ_2016-01-01_2025-12-31_yahoo_daily.csv"
    prior = pd.read_csv(prior_path, parse_dates=["date"]).sort_values("date")
    qqq = pd.read_csv(qqq_path, parse_dates=["date"]).sort_values("date")
    price_col = "adj_close" if "adj_close" in qqq.columns else ("Adj Close" if "Adj Close" in qqq.columns else "close")
    qqq["qqq_equity"] = initial_capital * qqq[price_col].astype(float) / float(qqq[price_col].iloc[0])
    qqq["qqq_dd"] = qqq["qqq_equity"] - qqq["qqq_equity"].cummax()
    prior["prior_dd"] = prior["stress_equity"] - prior["closed_equity"].cummax()

    annual = prior.groupby(prior["date"].dt.year)["daily_pnl"].sum()
    qqq_annual = qqq.groupby(qqq["date"].dt.year).agg(first=(price_col, "first"), last=(price_col, "last"))
    qqq_annual["return_pct"] = (qqq_annual["last"] / qqq_annual["first"] - 1.0) * 100.0
    covid_prior_return = float(annual.loc[2020] / initial_capital * 100.0) if 2020 in annual.index else float("nan")
    covid_qqq_return = float(qqq_annual.loc[2020, "return_pct"]) if 2020 in qqq_annual.index else float("nan")

    prior_end = float(prior["closed_equity"].iloc[-1])
    prior_net = prior_end - initial_capital
    prior_stress_dd = float(prior["stress_dd"].min())
    qqq_end = float(qqq["qqq_equity"].iloc[-1])
    qqq_net = qqq_end - initial_capital
    qqq_dd = float(qqq["qqq_dd"].min())
    rows = [
        {
            "Sleeve": "Prior-opposed intraday NQ fixed 1 base book",
            "End Capital": money(prior_end, 0),
            "Net": money(prior_net, 0),
            "Max DD": money(prior_stress_dd, 0),
            "Return": pct(prior_net / initial_capital * 100.0, 1),
            "Net/DD": ratio(prior_net / abs(prior_stress_dd), 2),
            "Peak Size": "5 units",
        },
        {
            "Sleeve": "QQQ fully invested",
            "End Capital": money(qqq_end, 0),
            "Net": money(qqq_net, 0),
            "Max DD": money(qqq_dd, 0),
            "Return": pct(qqq_net / initial_capital * 100.0, 1),
            "Net/DD": ratio(qqq_net / abs(qqq_dd), 2),
            "Peak Size": "full ETF capital",
        },
    ]

    chart = pd.DataFrame({"date": prior["date"], "prior_equity": prior["closed_equity"]}).merge(
        qqq[["date", "qqq_equity"]],
        on="date",
        how="outer",
    ).sort_values("date")
    facts = pd.Series(
        {
            "covid_prior_return_pct": covid_prior_return,
            "covid_qqq_return_pct": covid_qqq_return,
            "start": prior["date"].min(),
            "end": prior["date"].max(),
        }
    )
    return chart, rows, facts


def funding_benchmark_rows(summaries: Dict[str, Dict[str, object]]) -> List[Dict[str, object]]:
    """Two-system capital-efficiency view for the funding deck.

    The prior-opposed row is read from the top-strategy benchmark CSV when
    available. The ungated row is computed from the same stress-capital anchor
    so the two proposed systems can be shown side by side without bringing
    unrelated strategy families into the buyer-facing deck.
    """

    normalized_path = FAIR_BENCHMARK_DIR / "top_strats_max_stress_normalized.csv"
    executable_path = FAIR_BENCHMARK_DIR / "top_strats_common_account_executable.csv"
    normalized = pd.read_csv(normalized_path) if normalized_path.exists() else pd.DataFrame()
    executable = pd.read_csv(executable_path) if executable_path.exists() else pd.DataFrame()

    common_capital = float(normalized["start_capital"].iloc[0]) if not normalized.empty else 927_205.5
    common_stress = float(abs(normalized["scaled_stress_dd_usd"].iloc[0])) if not normalized.empty else 309_068.5
    common_account = float(executable["start_capital"].iloc[0]) if not executable.empty else 1_000_000.0

    def public_row(label: str, src_slug: str, normalized_row: Optional[pd.Series] = None, executable_row: Optional[pd.Series] = None) -> Dict[str, object]:
        s = summaries[src_slug]
        base_net = float(s["total_net"])
        base_stress = abs(float(s["max_stress_dd_usd"]))
        required_3x = base_stress * 3.0

        if normalized_row is not None:
            norm_scale = float(normalized_row["scale_factor_base_books"])
            norm_net = float(normalized_row["scaled_net_usd"])
            norm_return = float(normalized_row["scaled_return_pct"])
            norm_net_dd = float(normalized_row["scaled_net_over_stress"])
            norm_qqq_dca = float(normalized_row["qqq_dca_net_usd"])
        else:
            norm_scale = common_capital / required_3x if required_3x else 0.0
            norm_net = base_net * norm_scale
            norm_return = norm_net / common_capital * 100.0 if common_capital else 0.0
            norm_net_dd = norm_net / common_stress if common_stress else 0.0
            norm_qqq_dca = float(
                normalized.loc[normalized["window_start"].eq(str(s["start_date"].date())), "qqq_dca_net_usd"].iloc[0]
            ) if not normalized.empty and normalized["window_start"].eq(str(s["start_date"].date())).any() else float("nan")

        if executable_row is not None:
            books = int(executable_row["integer_books"])
            exec_net = float(executable_row["futures_net_usd"])
            exec_return = float(executable_row["futures_return_pct"])
            exec_net_dd = float(executable_row["futures_net_over_stress"])
            exec_qqq_dca = float(executable_row["qqq_dca_net_usd"])
        else:
            books = int(common_account // required_3x) if required_3x else 0
            exec_net = base_net * books
            exec_return = exec_net / common_account * 100.0 if common_account else 0.0
            exec_net_dd = exec_net / (base_stress * books) if base_stress and books else 0.0
            exec_qqq_dca = float(
                executable.loc[executable["window_start"].eq(str(s["start_date"].date())), "qqq_dca_net_usd"].iloc[0]
            ) if not executable.empty and executable["window_start"].eq(str(s["start_date"].date())).any() else float("nan")

        return {
            "System": label,
            "Window": f"{s['start_date'].date()} to {s['end_date'].date()}",
            "Base Net": money(base_net, 0),
            "Base Stress": money(-base_stress, 0),
            "Normalized Scale": f"{norm_scale:.2f}x",
            "Normalized Net": money(norm_net, 0),
            "Normalized Return": pct(norm_return, 1),
            "Norm. Net/DD": ratio(norm_net_dd, 2),
            "Norm. vs QQQ DCA": money(norm_net - norm_qqq_dca, 0) if np.isfinite(norm_qqq_dca) else "n/a",
            "$1M Books": str(books),
            "$1M Net": money(exec_net, 0),
            "$1M Return": pct(exec_return, 1),
            "$1M Net/DD": ratio(exec_net_dd, 2),
            "$1M vs QQQ DCA": money(exec_net - exec_qqq_dca, 0) if np.isfinite(exec_qqq_dca) else "n/a",
        }

    prior_norm = None
    if not normalized.empty:
        rows = normalized[normalized["slug"].eq("nq_v2b_prior_opposed_stpmc_only_S_1_1_3")]
        prior_norm = rows.iloc[0] if not rows.empty else None
    prior_exec = None
    if not executable.empty:
        rows = executable[executable["slug"].eq("nq_v2b_prior_opposed_stpmc_only_S_1_1_3")]
        prior_exec = rows.iloc[0] if not rows.empty else None

    return [
        public_row("Gated intraday", "nq_prior_opposed_full", prior_norm, prior_exec),
        public_row("Ungated intraday", "nq_ungated_intraday"),
    ]


def pitch_deck(summaries: Dict[str, Dict[str, object]]) -> str:
    petit = summaries["mnq_prior_opposed"]
    grand = summaries["nq_prior_opposed_full"]
    rows = summary_rows(summaries)
    benchmark_rows = funding_benchmark_rows(summaries)
    qqq_rows = []
    for slug in ["mnq_prior_opposed", "mnq_ungated_intraday", "nq_prior_opposed_full", "nq_ungated_intraday"]:
        s = summaries[slug]
        qqq_rows.append(
            {
                "System": f"{s['market']} {s['public_name']}",
                "Corr": ratio(float(s["daily_corr_to_qqq"]), 2),
                "Avg in QQQ Down Years": pct(float(s["avg_return_in_qqq_down_years_pct"]), 1),
                "QQQ Down Years": s["qqq_down_years"] or "none",
            }
        )

    return "\n".join(
        [
            f"# {PROGRAM_NAME} - Private Funding Deck",
            "",
            "**Status:** private diligence draft for family/friends funding discussions. All trading results are hypothetical/backtested, unaudited, and not a live managed-account track record.",
            "",
            "---",
            "",
            "## 1. The Ask",
            "",
            f"- Seeking **about {money(STARTER_ASK, 0)}** to fund a CTA-registration-first {RUNWAY_MONTHS}-month build/test runway.",
            "- Use of funds: platform development, live data, broker-paper validation, CTA registration/disclosure work, and segregated MNQ research/test capital.",
            "- This is not framed as an immediate public CTA solicitation or a promise of returns.",
            "- No compensated futures advice, managed client trading, or client account direction should begin until counsel confirms the required registration/exemption, disclosure, and account-documentation path.",
            "- The initial research and validation path stays on MNQ to keep operating risk and research cost controlled.",
            "",
            "---",
            "",
            "## 2. Why This Exists",
            "",
            "- Most prospective investors already have stock and real-estate exposure.",
            "- The product goal is a futures-based diversification tool that can behave differently from QQQ.",
            "- Existing research shows slightly negative daily return correlation to QQQ across the selected intraday systems.",
            "- The build phase is meant to prove execution, reporting, and risk controls before scaling.",
            "",
            "---",
            "",
            "## 3. Founder Partner Exchange",
            "",
            "Family and friends would be taking the earliest and most personal risk: funding an unproven live-test build before there is an audited managed-account record. The value of that support is not just capital; it creates the conditions for an auditable operating history, live execution evidence, and a future CTA-ready platform.",
            "",
            md_table(
                [
                    {
                        "What They Provide": "Early risk capital and patience",
                        "Proposed Exchange": "Priority access to the earliest validated capacity if the program launches",
                    },
                    {
                        "What They Provide": "Trust before an audited track record exists",
                        "Proposed Exchange": f"Founder economics target: {FOUNDER_PARTNER_MGMT_FEE:.0f}% management / {FOUNDER_PARTNER_PERFORMANCE_FEE:.0f}% performance fee, subject to counsel-approved documents",
                    },
                    {
                        "What They Provide": "A path to build a real operating record",
                        "Proposed Exchange": "Option to carry their validated allocation into a future compliant CTA-style product or withdraw after review windows",
                    },
                ],
                ["What They Provide", "Proposed Exchange"],
            ),
            "",
            "If the firm later scales, these accounts should be documented through counsel-reviewed founder terms or side letters. This is not a guarantee of returns, capacity, fee treatment, or a future offering; it is the intended alignment principle for people who help make the platform possible.",
            "",
            "---",
            "",
            "## 4. Product Tiers",
            "",
            md_table(
                [
                    {
                        "Tier": PETIT_NAME,
                        "Market": "MNQ",
                        "Minimum": money(MNQ_MINIMUM, 0),
                        "Use": "Lower-cost research/live-test tier",
                    },
                    {
                        "Tier": GRAND_NAME,
                        "Market": "NQ",
                        "Minimum": money(NQ_MINIMUM, 0),
                        "Use": "Future larger-account tier after validation",
                    },
                ],
                ["Tier", "Market", "Minimum", "Use"],
            ),
            "",
            "---",
            "",
            "## 5. Systems Under Development",
            "",
            "- **Prior-Opposed Gated Intraday System:** flagship conditional intraday strategy. It waits for a proprietary market condition before allowing a later opposing campaign.",
            "- **Ungated Intraday Breakout System:** secondary all-day intraday strategy. It is simpler and useful as a benchmark, but less selective.",
            "- Exact signal formulas, timing rules, and sizing maps are intentionally omitted from this private deck.",
            "",
            "---",
            "",
            "## 6. Conservative Performance Framing",
            "",
            "Headline expected-return framing removes each system's top 3 annual net years, then averages the remaining annual results.",
            "",
            md_table(rows, ["Tier", "System", "Market", "Minimum", "Window", "Total Net", "Advertised Avg/Yr", "Advertised Return/Yr", "Corr to QQQ"]),
            "",
            "---",
            "",
            "## 7. QQQ Diversification Evidence",
            "",
            "QQQ comparisons follow the fair-benchmark convention: the passive benchmark invests the full tier minimum in QQQ over the same clipped strategy window, while the futures row starts with that same tier minimum plus realized strategy P&L.",
            "",
            md_table(qqq_rows, ["System", "Corr", "Avg in QQQ Down Years", "QQQ Down Years"]),
            "",
            "![Daily correlation to QQQ](charts/daily_correlation_to_qqq.png)",
            "",
            "---",
            "",
            "## 8. Preferred Tier Equity vs QQQ",
            "",
            "![Preferred tiers vs QQQ](charts/preferred_tiers_vs_qqq_equity.png)",
            "",
            "---",
            "",
            "## 9. Capital-Efficiency Context For Test Systems",
            "",
            "This table keeps the buyer-facing view focused on the two systems planned for development. The normalized column uses the same common stress-capital anchor as the research tracker, while the executable column shows whole-book feasibility on a $1,000,000 account.",
            "",
            md_table(
                [
                    {
                        "System": row["System"],
                        "Norm Scale": row["Normalized Scale"],
                        "Norm Net": row["Normalized Net"],
                        "Norm Return": row["Normalized Return"],
                        "Norm Net/DD": row["Norm. Net/DD"],
                        "$1M Books": row["$1M Books"],
                        "$1M Net": row["$1M Net"],
                        "$1M Return": row["$1M Return"],
                    }
                    for row in benchmark_rows
                ],
                ["System", "Norm Scale", "Norm Net", "Norm Return", "Norm Net/DD", "$1M Books", "$1M Net", "$1M Return"],
            ),
            "",
            "---",
            "",
            "## 10. Nasdaq Intraday Mini Year-By-Year",
            "",
            "![Mini annual returns vs QQQ](charts/mini_annual_returns_vs_qqq.png)",
            "",
            "---",
            "",
            f"## 11. {RUNWAY_MONTHS}-Month Build Roadmap",
            "",
            "1. **Months 1-3:** CTA counsel kickoff, Series 3/NFA registration prep, source-data audit, CQG/StoneX adapter work, live-data shadow mode.",
            "2. **Months 4-6:** disclosure-document drafting/review, broker-paper routing, reconciliation, restart drills, emergency flatten drills, first readiness review.",
            "3. **Months 7-9:** extended MNQ funded-paper or small-live observation, ungated-system comparison, investor reporting packet.",
            "4. **Months 10-12:** regime review, QQQ comparison refresh, compliance/accounting package, registration/disclosure go/no-go, final MNQ/NQ tier decision.",
            "",
            "---",
            "",
            f"## 12. {RUNWAY_MONTHS}-Month Budget",
            "",
            f"- Total target: **{money(STARTER_ASK, 0)}**.",
            "- Largest line items: founder development stipend and segregated MNQ research/test capital.",
            "- The cost model now assumes StoneX for futures execution, using the supplied commission and account-fee schedule.",
            "- CTA registration, disclosure-document counsel, audit/accounting setup, and NFA response reserve are now first-class budget items before taking or managing outside trading capital.",
            "",
            "See [COST_AND_RUNWAY.md](COST_AND_RUNWAY.md).",
            "",
            "---",
            "",
            "## 13. Risk Positioning",
            "",
            "- Futures are leveraged and can lose more than expected when markets gap or systems fail.",
            "- Backtests are not live results.",
            "- The platform must prove feed integrity, order handling, reconciliation, and reporting before any client-scale allocations.",
            "- Materials and account documents need counsel review before being used for solicitation or advisory activity.",
            "- A human kill switch, daily flattening, stale-feed blocks, and broker/local position reconciliation are required controls.",
            "",
            "---",
            "",
            "## 14. Next Step",
            "",
            f"- Fund the CTA-registration-first {RUNWAY_MONTHS}-month build/test runway.",
            "- Review monthly operating reports.",
            "- Decide after the pilot and counsel/regulatory review whether the platform is ready for broader family/friends capital.",
            "",
            "Internal data exhibits live in `exhibits/`. Public-facing numbers should be reviewed by counsel before distribution.",
        ]
    )


def short_memo(summaries: Dict[str, Dict[str, object]]) -> str:
    rows = summary_rows(summaries, preferred_only=True)
    petit_rows = annual_rows_for("mnq_prior_opposed")
    grand_rows = annual_rows_for("nq_prior_opposed_full")
    return "\n".join(
        [
            f"# {PROGRAM_NAME} - Short Funding Memo",
            "",
            f"**Purpose:** request funding for a CTA-registration-first {RUNWAY_MONTHS}-month development and validation runway. This memo is not legal, tax, or investment advice.",
            "",
            "## Executive Summary",
            "",
            f"The proposed project is a private build-and-test effort for a futures-based diversification platform. The first deployable tier is **{PETIT_NAME}**, using MNQ so research, live data, and operational validation can be done with lower capital intensity. **{GRAND_NAME}** is the future NQ version once the MNQ runtime proves stable.",
            "",
            f"The requested initial budget is **about {money(STARTER_ASK, 0)}**. That covers development time, data, cloud/runtime costs, broker/API/exchange allowances, CTA registration/disclosure work, audit/accounting reserve, contingency, and **{money(MNQ_MINIMUM, 0)}** of segregated MNQ research/test capital.",
            "",
            "## CTA Registration-First Adjustment",
            "",
            "The package now assumes the regulatory path is handled before any compensated commodity-interest advisory activity or managed client trading begins. The budget includes NFA/Series 3/registration administration and a larger CTA disclosure/legal/compliance/audit reserve. Canadian CTC/cross-border registration remains a separate counsel-confirmed expansion item, not a base-runway promise.",
            "",
            "## Founder Partner Exchange",
            "",
            "The early family/friends risk is real: they are supporting an unproven live-test platform with no audited managed-account track record. In exchange, the intended alignment is that successful early supporters become **Founder Partner Accounts** if the platform later transitions into a compliant CTA-style offering.",
            "",
            md_table(
                [
                    {
                        "Founder Benefit": "Early access",
                        "Draft Principle": "They can participate from day zero in a systematic futures program if the live validation supports launch.",
                    },
                    {
                        "Founder Benefit": "Better economics",
                        "Draft Principle": f"Target founder terms are {FOUNDER_PARTNER_MGMT_FEE:.0f}% management / {FOUNDER_PARTNER_PERFORMANCE_FEE:.0f}% performance fee instead of a future standard 1/20 model.",
                    },
                    {
                        "Founder Benefit": "Capacity priority",
                        "Draft Principle": "If the strategy later becomes capacity-constrained, early accounts should receive priority access before new outside allocations.",
                    },
                    {
                        "Founder Benefit": "Formal protection",
                        "Draft Principle": "Any founder terms should be documented in counsel-reviewed side letters or equivalent account documents.",
                    },
                ],
                ["Founder Benefit", "Draft Principle"],
            ),
            "",
            "These are proposed alignment principles, not a guarantee of returns, fee treatment, capacity, or future offering availability. Counsel must decide the correct documentation before any investor-facing commitment is made.",
            "",
            "## Why It May Diversify QQQ Exposure",
            "",
            "The selected futures strategies are intraday and can trade both directions. They are not designed as passive Nasdaq exposure. QQQ benchmark rows follow the same fair-comparison convention used in the research tracker: the passive benchmark invests the full tier minimum over the same clipped strategy window. In the current backtests, daily return correlations to QQQ are slightly negative, and the preferred systems remained positive during the available QQQ down-year windows.",
            "",
            md_table(rows, ["Tier", "System", "Market", "Minimum", "Window", "Total Net", "Advertised Avg/Yr", "Advertised Return/Yr", "Corr to QQQ"]),
            "",
            "## Conservative Return Standard",
            "",
            "For conversations with prospective investors, the headline number should not use the full average annual return. The package removes the top 3 annual net years for each system and averages only the remaining years. Full annual tables remain visible below so the good, weak, and partial years are all represented.",
            "",
            "## Nasdaq Intraday Mini Annual Table",
            "",
            md_table(petit_rows, ["Year", "System Return", "Net", "QQQ Return", "QQQ Regime"]),
            "",
            "## Nasdaq Intraday Annual Table",
            "",
            md_table(grand_rows, ["Year", "System Return", "Net", "QQQ Return", "QQQ Regime"]),
            "",
            f"## What The {RUNWAY_MONTHS} Months Must Prove",
            "",
            "- The platform can ingest live data, form bars, and reproduce replay signals after market close.",
            "- Broker-paper orders match local order intent, fills, and positions.",
            "- Daily flattening, stale-feed blocks, and kill-switch behavior work under real operating conditions.",
            "- Monthly reports can show equity, drawdown, exposure, fills, costs, incidents, and reconciliation status.",
            "- The strategy still makes operational sense after live shadow/paper evidence, not only historical replay.",
            "",
            "## What Is Out Of Scope For The First Ask",
            "",
            "- No promise of a live CTA track record.",
            "- No public solicitation.",
            "- No compensated CTA advice, client account direction, or managed futures trading until counsel confirms registration/exemption and disclosure requirements are satisfied.",
            "- No NQ live trading before MNQ execution is validated.",
            "- No guarantee that historical returns will persist.",
            "",
            "## Diligence Notes",
            "",
            "The strategy mechanics are intentionally summarized at a high level. Exact formulas, state gates, and sizing maps should stay internal or be disclosed only under a separate diligence process. Any investor-facing version needs counsel review before distribution.",
        ]
    )


def cost_and_runway() -> str:
    budget = pd.read_csv(EXHIBITS / BUDGET_EXHIBIT)
    rows = [
        {"Category": r.Category, "Basis": r.Basis, "Amount": money(float(r.Amount), 0)}
        for r in budget.itertuples(index=False)
    ]
    regulatory = pd.read_csv(EXHIBITS / "regulatory_cost_assumptions.csv")
    regulatory_rows = [
        {"Item": r.Item, "Assumption": r.Assumption, "Source": r.Source}
        for r in regulatory.itertuples(index=False)
    ]
    stonex = pd.read_csv(EXHIBITS / "stonex_fee_estimate.csv")
    fee_rows = []
    for row in stonex.itertuples(index=False):
        fee_rows.append(
            {
                "System": f"{row.market} {row.strategy}",
                "Closed Units": f"{int(row.closed_units):,}",
                "StoneX Commission": money(float(row.stonex_commission_usd), 0),
                "Model Audit Fee": money(float(row.model_audit_fee_usd), 0),
                "Delta": money(float(row.stonex_delta_vs_model_fee_usd), 0),
            }
        )
    stonex_schedule_rows = [
        {"Item": "Retail brokerage account minimum", "Provided Schedule": "$2,000 baseline"},
        {"Item": "Live simulation account minimum", "Provided Schedule": "$5,000 baseline"},
        {"Item": "Standard futures contracts", "Provided Schedule": "$1.29 per side, per contract"},
        {"Item": "Micro futures contracts", "Provided Schedule": "$0.50 per side, per contract"},
        {"Item": "StoneX Futures Platform", "Provided Schedule": "Free"},
        {"Item": "Mobile app trading", "Provided Schedule": "Free"},
        {"Item": "Simulation software", "Provided Schedule": "$49.95/month"},
        {"Item": "Domestic / international outbound wire", "Provided Schedule": "$25 / $50"},
        {"Item": "ACAT transfer out", "Provided Schedule": "$125"},
        {"Item": "Inactive account fee", "Provided Schedule": "$100/year"},
        {"Item": "Margin debit rate", "Provided Schedule": "WSJ Call Money Rate + 2.50%"},
    ]
    return "\n".join(
        [
            "# Cost And Runway Plan",
            "",
            f"Target raise: **about {money(STARTER_ASK, 0)}** for a CTA-registration-first {RUNWAY_MONTHS}-month build/test runway.",
            "",
            md_table(rows, ["Category", "Basis", "Amount"]),
            "",
            f"Planned budget total: **{money(float(budget['Amount'].sum()), 0)}**.",
            "",
            "## Regulatory Cost Assumptions",
            "",
            "The registration budget separates official NFA-style fees from planning reserves. The disclosure/legal/compliance/audit reserve is intentionally larger than the hard filing fees because a usable CTA package needs counsel review, disclosure drafting, NFA response time, recordkeeping setup, and accounting/audit structure before any client-facing trading activity.",
            "",
            md_table(regulatory_rows, ["Item", "Assumption", "Source"]),
            "",
            "## StoneX Fee Schedule Assumption",
            "",
            "The following table uses the StoneX One fee schedule supplied by the founder for this draft. It has not been independently verified inside a live StoneX account, so counsel/account-opening review should confirm it before funding.",
            "",
            md_table(stonex_schedule_rows, ["Item", "Provided Schedule"]),
            "",
            "## StoneX Commission Estimate For Current Models",
            "",
            f"The replay engine already embeds a **{money(MODEL_AUDIT_FEE_PER_CLOSED_UNIT, 2)} per closed-unit audit fee**. The table below compares that embedded audit fee to StoneX commission-only estimates using the supplied per-side rates. Exchange, NFA, clearing, market-data, margin-interest, wire, and slippage costs are not included in these commission-only rows.",
            "",
            md_table(fee_rows, ["System", "Closed Units", "StoneX Commission", "Model Audit Fee", "Delta"]),
            "",
            "## Runway Milestones",
            "",
            "| Month | Milestone | Required Evidence |",
            "| --- | --- | --- |",
            "| 1 | CTA counsel kickoff and runtime hardening | Structure memo, source manifest, reproducible research package |",
            "| 2 | Series 3 / NFA ORS readiness and live-data shadow mode | Exam/admin plan, stored live bars, no-trade signal reports |",
            "| 3 | Disclosure document draft and broker-paper adapter | Draft disclosure outline, broker order ids mapped to local intents |",
            "| 4 | MNQ broker-paper trial | Daily reports, slippage/cost audit, incident log |",
            "| 5 | CQG/StoneX demo hardening and NFA response reserve | Account/order reconciliation, restart drill, emergency flatten drill |",
            "| 6 | First readiness review | Feed integrity report, order sequencing audit, regulatory gap list |",
            "| 7 | Extended MNQ funded-paper or small-live continuation | Stable reports, variance-to-replay audit, risk-limit adherence |",
            "| 8 | Ungated v2b paper comparison | Secondary system replay-vs-paper evidence and operational differences |",
            "| 9 | Reporting and investor portal draft | Monthly packet, exposure report, drawdown explanation template |",
            "| 10 | Robustness and regime review | QQQ comparison refresh, bad-market behavior, filter review |",
            "| 11 | Compliance/accounting package | Counsel checklist, hypothetical-performance language, recordkeeping plan |",
            "| 12 | Final go/no-go and tier decision | Pilot report, registration/disclosure status, NQ tier decision only after MNQ evidence |",
            "",
            "## Capital Treatment",
            "",
            f"- The **{money(MNQ_MINIMUM, 0)} MNQ research/test capital** should be segregated from development spend.",
            "- The NQ tier is not funded in this first runway; it remains a future tier after MNQ operations are stable.",
            "- The MNQ test capital is internal research/test capital unless counsel confirms a compliant managed-account or advisory structure.",
            "- Any client trading funds require CTA registration or a confirmed exemption, accepted disclosure/account documents if required, and broker/account approvals before acceptance.",
            "- Canadian CTC/cross-border registration is not included in the base raise; it should be separately scoped if Canadian client activity is pursued.",
            "",
            "## Operating Cost Notes",
            "",
            "- Databento live data is estimated from the current internal spec at $179/month; commercial/non-display classification can raise the actual cost.",
            "- StoneX account minimums and fee schedule should be confirmed inside the account-opening paperwork before live use.",
            "- Broker/API/exchange fees must be confirmed inside the broker and data-provider portals before live use.",
            "- Development stipend is included because the platform requires real engineering work: data ingest, order routing, reconciliation, reporting, deployment, monitoring, and documentation.",
        ]
    )


def reporting_templates() -> str:
    return "\n".join(
        [
            "# Reporting Templates",
            "",
            "These templates are inspired by the sample CTA equity and position workbooks in `data/CTA Samples/`. They are written as Markdown shapes first; they can later become CSV/XLSX/PDF reports.",
            "",
            "## Monthly Investor Summary",
            "",
            "| Field | Value |",
            "| --- | --- |",
            "| Account / Tier |  |",
            "| Beginning Equity |  |",
            "| Ending Equity |  |",
            "| Net P&L |  |",
            "| Monthly Return |  |",
            "| Year-To-Date Return |  |",
            "| Max Month Drawdown |  |",
            "| Intrabar Stress Drawdown |  |",
            "| Fees / Commissions / Data Allocations |  |",
            "| Open Positions At Month End |  |",
            "| Operational Incidents |  |",
            "| Continue / Pause / Review Decision |  |",
            "",
            "## Equity And Margin Status",
            "",
            "| Date | Account | Tier | Start Balance | Net Liquidation Value | Initial Margin | Margin Excess | Percent Margin Excess | Notes |",
            "| --- | --- | --- | ---: | ---: | ---: | ---: | ---: | --- |",
            "| YYYY-MM-DD | Example | Mini |  |  |  |  |  |  |",
            "",
            "## Daily Equity Summary",
            "",
            "| Date | Beginning Equity | Ending Equity | Day P&L | Day Return | Month Return | Year Return | Drawdown | Reconciled? |",
            "| --- | ---: | ---: | ---: | ---: | ---: | ---: | ---: | --- |",
            "| YYYY-MM-DD |  |  |  |  |  |  |  |  |",
            "",
            "## Position / Exposure Summary",
            "",
            "| Date | Strategy | Instrument | Direction | Contracts | Entry Time | Exit Time | Realized P&L | Broker Order IDs Present? | Local/Broker Position Match? |",
            "| --- | --- | --- | --- | ---: | --- | --- | ---: | --- | --- |",
            "| YYYY-MM-DD |  |  |  |  |  |  |  |  |  |",
            "",
            "## Execution Fidelity Report",
            "",
            "| Check | Pass/Fail | Count | Notes |",
            "| --- | --- | ---: | --- |",
            "| Live bars persisted |  |  |  |",
            "| End-of-day replay matched live signals |  |  |  |",
            "| Broker order ids mapped to local intents |  |  |  |",
            "| Broker/local position mismatches |  |  |  |",
            "| Missed EOD flatten events |  |  |  |",
            "| Stale-feed entry blocks |  |  |  |",
            "| Manual interventions |  |  |  |",
            "| Unexpected orders/fills |  |  |  |",
            "",
            "## Monthly Cost Report",
            "",
            "| Cost Type | Amount | Notes |",
            "| --- | ---: | --- |",
            "| Market data |  |  |",
            "| Broker/API/exchange fees |  |  |",
            "| Commissions |  |  |",
            "| Slippage estimate |  |  |",
            "| Cloud/runtime |  |  |",
            "| Legal/compliance/accounting |  |  |",
        ]
    )


def risk_and_disclosure_notes() -> str:
    return "\n".join(
        [
            "# Risk And Disclosure Notes",
            "",
            "This is a working checklist, not legal advice. A qualified attorney/compliance consultant should review any material before funds are accepted or performance is distributed.",
            "",
            "## Performance Status",
            "",
            "- All trading results in this package are hypothetical/backtested.",
            "- No audited live CTA track record is presented.",
            "- Results include model assumptions, replay assumptions, and historical data limitations.",
            "- Hypothetical results may differ materially from actual trading due to liquidity, market impact, slippage, missed orders, platform outages, human behavior, and changing market regimes.",
            "",
            "## Futures Risk",
            "",
            "- Futures are leveraged and can lose more than the posted margin or planned allocation.",
            "- Stops may fill worse than expected during gaps or fast markets.",
            "- Technology failure can leave orders, positions, or data state out of sync.",
            "- The platform must include stale-feed controls, broker/local reconciliation, daily flattening, and a manual kill switch.",
            "- StoneX fees in this package are based on the founder-provided schedule and should be verified during account opening.",
            "",
            "## Private Offering / CTA Review Checklist",
            "",
            "- Determine whether the structure is an advisory relationship, pooled vehicle, software/service subscription, or another legal form.",
            "- Confirm whether CTA registration or an exemption applies before managing or advising commodity-interest accounts.",
            "- Treat CTA registration as the default working path unless counsel documents a valid exemption for the actual fundraising and advisory structure.",
            "- Do not provide compensated futures advice, place discretionary trades, or direct client commodity-interest accounts until registration/exemption and disclosure requirements are satisfied.",
            "- Confirm whether any securities offering exemption is needed if raising operating capital from family/friends.",
            "- Avoid public solicitation unless counsel confirms the structure permits it.",
            "- Verify investor suitability/accredited status requirements where applicable.",
            "- Keep hypothetical-performance disclaimers next to performance tables, not only at the end.",
            "- Maintain support files for every performance claim, including source data, replay assumptions, commissions/fees, and annual return calculations.",
            "",
            "## CTA Registration-First Checklist",
            "",
            "- Designate an NFA ORS Security Manager and prepare online Form 7-R.",
            "- Identify principals/APs, Form 8-R requirements, fingerprints, and proficiency requirements such as Series 3.",
            "- Budget for the NFA application fee, CTA dues, AP/principal fee, exam/admin costs, counsel, disclosure drafting, and accounting/audit setup.",
            "- Prepare the CTA disclosure document and do not use it with prospective clients until required NFA review/acceptance is complete.",
            "- Review all hypothetical-performance and promotional materials against NFA Rule 2-29 before distribution.",
            "- Confirm whether Canadian CTC/cross-border registration, NRD filings, or Canadian working-capital rules apply before taking Canadian client activity.",
            "",
            "## Founder Partner / Side Letter Checklist",
            "",
            "- Treat founder economics as a draft alignment principle until counsel documents them.",
            f"- Proposed founder economics in this package are {FOUNDER_PARTNER_MGMT_FEE:.0f}% management / {FOUNDER_PARTNER_PERFORMANCE_FEE:.0f}% performance fee, compared with a possible future standard 1/20 model.",
            "- Do not promise permanent capacity, permanent fee terms, or withdrawal rights unless they are written into counsel-approved agreements.",
            "- Confirm whether special founder terms create conflicts, disclosure obligations, side-letter parity requirements, or most-favored-nation issues.",
            "- If family/friends later transition into managed accounts, make sure every account's fee treatment and high-water-mark mechanics are auditable.",
            "",
            "## Official References",
            "",
            "- SEC accredited investor guidance: https://www.sec.gov/resources-small-businesses/capital-raising-building-blocks/accredited-investors",
            "- Investor.gov Rule 506 overview: https://www.investor.gov/introduction-investing/investing-basics/glossary/rule-506-regulation-d",
            "- NFA CTA registration requirements: https://www.nfa.futures.org/registration-membership/who-has-to-register/cta.html",
            "- NFA membership dues and fees: https://www.nfa.futures.org/registration-membership/dues-revenue-structure.html",
            "- NFA CTA disclosure documents: https://www.nfa.futures.org/members/cta/regulatory-obligations/disclosure-documents.html",
            "- NFA CTA disclosure document FAQs: https://www.nfa.futures.org/faqs/members/cta-disclosure-documents.html",
            "- NFA Compliance Rule 2-29: https://www.nfa.futures.org/rulebooksql/rules.aspx?RuleID=RULE+2-29&Section=4",
            "",
            "## Language To Avoid Before Counsel Review",
            "",
            "- Do not say results are guaranteed, expected, safe, proven, or audited.",
            "- Do not imply the backtest is equivalent to actual client trading.",
            "- Do not advertise exact proprietary formulas in general materials.",
            "- Do not present the NQ tier as live-ready before MNQ operations are validated.",
        ]
    )


def write_powerpoint(summaries: Dict[str, Dict[str, object]]) -> None:
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]

    navy = RGBColor(17, 36, 58)
    blue = RGBColor(45, 100, 165)
    green = RGBColor(36, 139, 96)
    gold = RGBColor(206, 151, 48)
    slate = RGBColor(84, 96, 112)
    pale = RGBColor(248, 250, 252)
    white = RGBColor(255, 255, 255)
    dark = RGBColor(20, 28, 38)
    border = RGBColor(219, 226, 235)

    def set_bg(slide, color=pale):
        fill = slide.background.fill
        fill.solid()
        fill.fore_color.rgb = color

    def add_box(slide, x, y, w, h, fill_color=white, line_color=border, radius=True):
        shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE
        shape = slide.shapes.add_shape(shape_type, Inches(x), Inches(y), Inches(w), Inches(h))
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
        shape.line.color.rgb = line_color
        shape.line.width = Pt(0.75)
        return shape

    def add_text(
        slide,
        text,
        x,
        y,
        w,
        h,
        size=18,
        bold=False,
        color=dark,
        align=PP_ALIGN.LEFT,
        font_name="Aptos",
    ):
        shape = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
        tf = shape.text_frame
        tf.clear()
        tf.word_wrap = True
        tf.margin_left = Inches(0.02)
        tf.margin_right = Inches(0.02)
        tf.margin_top = Inches(0.02)
        tf.margin_bottom = Inches(0.02)
        tf.vertical_anchor = MSO_ANCHOR.TOP
        p = tf.paragraphs[0]
        p.alignment = align
        r = p.add_run()
        r.text = text
        r.font.name = font_name
        r.font.size = Pt(size)
        r.font.bold = bold
        r.font.color.rgb = color
        return shape

    def add_bullets(slide, bullets, x, y, w, h, size=16, color=dark):
        shape = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
        tf = shape.text_frame
        tf.clear()
        tf.word_wrap = True
        tf.margin_left = Inches(0.02)
        tf.margin_right = Inches(0.02)
        tf.margin_top = Inches(0.02)
        for idx, bullet in enumerate(bullets):
            p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
            p.text = f"• {bullet}"
            p.font.name = "Aptos"
            p.font.size = Pt(size)
            p.font.color.rgb = color
            p.space_after = Pt(5)
        return shape

    def add_header(slide, title, kicker=None):
        set_bg(slide, white)
        add_text(slide, title, 0.62, 0.52, 9.4, 0.42, size=25, bold=False, color=blue)
        rule = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.64), Inches(0.98), Inches(8.85), Inches(0.012))
        rule.fill.solid()
        rule.fill.fore_color.rgb = RGBColor(190, 197, 205)
        rule.line.fill.background()
        if kicker:
            add_text(slide, kicker.upper(), 10.0, 0.61, 2.7, 0.22, size=8.5, color=RGBColor(170, 174, 181), align=PP_ALIGN.RIGHT)
        add_text(
            slide,
            "CONFIDENTIAL RESEARCH DRAFT - HYPOTHETICAL/BACKTESTED PERFORMANCE - COUNSEL REVIEW REQUIRED",
            2.15,
            7.2,
            8.9,
            0.18,
            size=6.5,
            color=RGBColor(137, 139, 144),
            align=PP_ALIGN.CENTER,
        )

    def add_metric_card(slide, label, value, x, y, w=2.8, h=1.0, accent=blue):
        add_box(slide, x, y, w, h)
        accent_shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(0.08), Inches(h))
        accent_shape.fill.solid()
        accent_shape.fill.fore_color.rgb = accent
        accent_shape.line.fill.background()
        add_text(slide, label, x + 0.22, y + 0.16, w - 0.35, 0.22, size=8.5, color=slate)
        add_text(slide, value, x + 0.22, y + 0.43, w - 0.35, 0.36, size=18, bold=True, color=dark)

    def add_table(slide, rows, columns, x, y, w, h, font_size=8.5, header_fill=navy, widths=None):
        shape = slide.shapes.add_table(len(rows) + 1, len(columns), Inches(x), Inches(y), Inches(w), Inches(h))
        table = shape.table
        if widths:
            for idx, width in enumerate(widths):
                table.columns[idx].width = Inches(width)
        for col_idx, col in enumerate(columns):
            cell = table.cell(0, col_idx)
            cell.text = col
            cell.fill.solid()
            cell.fill.fore_color.rgb = header_fill
            for p in cell.text_frame.paragraphs:
                for r in p.runs:
                    r.font.name = "Aptos"
                    r.font.size = Pt(font_size)
                    r.font.bold = True
                    r.font.color.rgb = white
        for row_idx, row in enumerate(rows, start=1):
            for col_idx, col in enumerate(columns):
                cell = table.cell(row_idx, col_idx)
                cell.text = str(row.get(col, ""))
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(255, 255, 255) if row_idx % 2 else RGBColor(241, 245, 249)
                for p in cell.text_frame.paragraphs:
                    p.alignment = PP_ALIGN.LEFT
                    for r in p.runs:
                        r.font.name = "Aptos"
                        r.font.size = Pt(font_size)
                        r.font.color.rgb = dark
        return shape

    def add_image_fit(slide, path, x, y, w, h):
        image_path = Path(path)
        with Image.open(image_path) as img:
            aspect = img.width / img.height
        box_aspect = w / h
        if aspect >= box_aspect:
            pic_w = w
            pic_h = w / aspect
        else:
            pic_h = h
            pic_w = h * aspect
        slide.shapes.add_picture(
            str(image_path),
            Inches(x + (w - pic_w) / 2),
            Inches(y + (h - pic_h) / 2),
            width=Inches(pic_w),
            height=Inches(pic_h),
        )

    scaling_daily, scaling_rows, scaling_facts = scaling_10y_nq_vs_qqq()

    # 1. Title
    slide = prs.slides.add_slide(blank)
    set_bg(slide, white)
    add_text(slide, PITCH_PROGRAM_NAME, 1.05, 2.78, 10.2, 0.6, size=34, bold=False, color=blue)
    title_rule = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1.06), Inches(3.4), Inches(9.95), Inches(0.012))
    title_rule.fill.solid()
    title_rule.fill.fore_color.rgb = RGBColor(181, 187, 194)
    title_rule.line.fill.background()
    add_text(slide, "INTRODUCTION", 1.15, 3.62, 4.8, 0.35, size=18, bold=False, color=RGBColor(172, 172, 176))
    add_text(slide, f"Private funding deck for a CTA-registration-first {RUNWAY_MONTHS}-month build/test runway", 1.16, 4.05, 7.95, 0.28, size=12.5, color=slate)
    add_text(slide, f"Target raise: {money(STARTER_ASK, 0)}  |  First test tier: {PETIT_NAME}", 1.16, 4.43, 7.5, 0.28, size=12.5, color=slate)
    add_text(
        slide,
        "CONFIDENTIAL RESEARCH DRAFT. HYPOTHETICAL/BACKTESTED PERFORMANCE IS NOT INDICATIVE OF FUTURE RESULTS.",
        1.05,
        7.12,
        10.9,
        0.18,
        size=7.3,
        color=RGBColor(137, 139, 144),
        align=PP_ALIGN.CENTER,
    )

    # 2. Research risk disclosure
    slide = prs.slides.add_slide(blank)
    add_header(slide, "Risk Disclosure Statement", "research draft")
    add_text(
        slide,
        "The risk of loss in futures trading can be substantial. Futures are leveraged instruments; adverse moves, gaps, liquidity events, data failures, or order-routing failures can produce losses materially different from historical tests.",
        0.85,
        1.35,
        11.7,
        0.78,
        size=13.5,
        color=dark,
    )
    add_text(
        slide,
        "This deck presents hypothetical and backtested research, not audited live CTA performance. It is intended to support a private build-and-validation discussion. No compensated futures advice, managed client trading, or client account direction should begin until counsel confirms the required registration/exemption, disclosure, broker, tax, accounting, and compliance path.",
        0.85,
        2.32,
        11.7,
        0.78,
        size=13.5,
        color=dark,
    )
    risk_rows = [
        {"Risk Area": "Execution", "Control Intent": "Match live order sequencing to post-session replay and broker-paper reviews."},
        {"Risk Area": "Sizing", "Control Intent": "Use historical intrabar stress drawdown as the account-sizing anchor."},
        {"Risk Area": "Operational", "Control Intent": "Stale-feed blocks, daily flattening, position reconciliation, and manual kill switch."},
        {"Risk Area": "Investor use", "Control Intent": "Private diligence only until CTA registration/exemption and counsel-approved materials exist."},
    ]
    add_table(slide, risk_rows, ["Risk Area", "Control Intent"], 0.85, 3.55, 11.55, 1.75, font_size=10.5, widths=[2.05, 9.5])
    add_text(slide, "Stops and protective orders cannot guarantee fills at intended prices during fast markets.", 0.9, 5.88, 10.5, 0.26, size=12.5, bold=True, color=blue)

    # 3. Fair Benchmark Scaling 10Y
    slide = prs.slides.add_slide(blank)
    add_header(slide, "Fair Benchmark Scaling 10Y", "$50k equity growth")
    add_image_fit(slide, CHARTS / "scaling_10y_nq_prior_vs_qqq_50k.png", 0.55, 1.2, 7.85, 3.9)
    sidebar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(8.65), Inches(0), Inches(4.68), Inches(7.5))
    sidebar.fill.solid()
    sidebar.fill.fore_color.rgb = blue
    sidebar.line.fill.background()
    add_text(slide, "Fixed one-base-book NQ study", 8.95, 1.25, 3.85, 0.34, size=18, bold=True, color=white)
    add_text(
        slide,
        "A $50k starting account is used here only as a fair benchmark scaling exhibit against fully invested QQQ. This fixed-base NQ study is not the 3x-stress live sizing recommendation.",
        8.95,
        1.82,
        3.85,
        1.15,
        size=12.2,
        color=white,
    )
    add_text(
        slide,
        "The exhibit shows the historical path that produced the table below. It is a research comparison, not a promise of capacity or future returns.",
        8.95,
        3.25,
        3.85,
        1.0,
        size=12.2,
        color=white,
    )
    add_text(slide, "2016-2025", 8.95, 5.1, 2.0, 0.3, size=20, bold=True, color=white)
    add_table(
        slide,
        scaling_rows,
        ["Sleeve", "End Capital", "Net", "Max DD", "Return", "Net/DD", "Peak Size"],
        0.45,
        5.55,
        7.95,
        0.92,
        font_size=7.4,
        widths=[2.0, 1.05, 1.05, 0.95, 0.85, 0.65, 1.4],
    )

    # 4. Program description
    slide = prs.slides.add_slide(blank)
    add_header(slide, "Program Description", "mechanical intraday futures")
    add_text(slide, "Nasdaq futures", 5.15, 1.4, 2.6, 0.34, size=19, bold=True, color=blue, align=PP_ALIGN.CENTER)
    desc_boxes = [
        ("Liquidity run", "The model observes an early directional move and waits for the market to reveal an adversarial setup rather than chasing every break."),
        ("Prior opposing condition", "A higher-timeframe proprietary condition, including prior-month reference context, must be known before the campaign can arm."),
        ("Mechanical price-action campaign", "Once eligible, execution is handled by a fixed intraday playbook with predefined entries, exits, stops, and session-end flattening."),
    ]
    xs = [0.8, 4.58, 8.36]
    for (title, body), x in zip(desc_boxes, xs):
        add_box(slide, x, 2.0, 3.25, 2.45, fill_color=RGBColor(252, 253, 255), radius=False)
        add_text(slide, title, x + 0.22, 2.28, 2.8, 0.28, size=16, bold=True, color=blue, align=PP_ALIGN.CENTER)
        add_text(slide, body, x + 0.28, 2.82, 2.7, 1.25, size=11.1, color=dark, align=PP_ALIGN.CENTER)
    add_box(slide, 2.08, 5.0, 9.18, 0.82, fill_color=RGBColor(246, 250, 253), radius=False)
    add_text(
        slide,
        "Optimization work remains conservative: reduce or downshift exposure when very large opening ranges historically degraded edge, while preserving most upside when the setup remains clean.",
        2.35,
        5.22,
        8.6,
        0.34,
        size=12.4,
        color=dark,
        align=PP_ALIGN.CENTER,
    )

    # 5. Main features
    slide = prs.slides.add_slide(blank)
    add_header(slide, "Main Features", "what the program is designed to do")
    features = [
        ("Agnostic but selective", "The strategy can accept either long or short exposure on a given day, but only after the proprietary condition allows it."),
        ("Adversarial NQ behavior", "The research is designed around the Nasdaq's tendency to punish the first obvious impulse and reward cleaner conditional follow-through."),
        ("Low market-beta dependence", "Daily returns in the current research set are not simply a proxy for QQQ. The strategy is meant to diversify stock-heavy portfolios."),
        ("Execution first", "The build phase focuses on matching live order flow, saved data, and post-session replay before any broader capital discussion."),
    ]
    y = 1.38
    for title, body in features:
        add_box(slide, 0.7, y, 3.15, 0.62, fill_color=white, line_color=blue, radius=False)
        add_text(slide, title, 0.9, y + 0.17, 2.72, 0.23, size=13.6, bold=True, color=blue, align=PP_ALIGN.CENTER)
        arrow = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(4.12), Inches(y + 0.18), Inches(0.55), Inches(0.24))
        arrow.fill.solid()
        arrow.fill.fore_color.rgb = blue
        arrow.line.fill.background()
        add_text(slide, body, 5.0, y + 0.02, 7.2, 0.52, size=12.5, color=dark)
        y += 1.22

    # 6. Strategy behavior in market regimes
    slide = prs.slides.add_slide(blank)
    add_header(slide, "Behavior In Different Market Regimes", "QQQ context")
    add_image_fit(slide, CHARTS / "daily_correlation_to_qqq.png", 0.7, 1.18, 5.65, 4.05)
    regime_rows = []
    for slug in ["mnq_prior_opposed", "nq_prior_opposed_full", "mnq_ungated_intraday", "nq_ungated_intraday"]:
        s = summaries[slug]
        regime_rows.append(
            {
                "System": f"{s['market']} {s['public_name'].split()[0]}",
                "Corr": ratio(float(s["daily_corr_to_qqq"]), 2),
                "QQQ Down-Year Avg": pct(float(s["avg_return_in_qqq_down_years_pct"]), 1),
            }
        )
    add_table(slide, regime_rows, ["System", "Corr", "QQQ Down-Year Avg"], 6.75, 1.25, 5.55, 1.58, font_size=9.2, widths=[2.25, 0.8, 2.5])
    add_bullets(
        slide,
        [
            "The strategy is not designed to be continuously long Nasdaq exposure.",
            f"In the 2020 shock/recovery year, the fixed NQ study returned {pct(float(scaling_facts['covid_prior_return_pct']), 1)} vs QQQ at {pct(float(scaling_facts['covid_qqq_return_pct']), 1)}.",
            "The weakest historical stress came from execution/range conditions rather than broad market direction alone.",
            "Large-range degradation filters reduce exposure in lower-quality regimes while trying to preserve the core upside.",
        ],
        6.95,
        3.28,
        5.0,
        1.9,
        size=11.7,
    )

    # 7. Source of edge and risk management
    slide = prs.slides.add_slide(blank)
    add_header(slide, "Source Of Edge & Risk Management", "execution quality before scale")
    add_text(slide, "Source of Edge", 0.78, 1.34, 3.2, 0.32, size=16, bold=True, color=blue)
    add_text(
        slide,
        "The research seeks to exploit sessions where a visible directional impulse creates an opposing, mechanically confirmable opportunity. The initial condition is not the main profit center; the follow-on campaign is.",
        0.78,
        1.82,
        11.55,
        0.78,
        size=12.7,
        color=dark,
    )
    add_text(slide, "Risk Management", 0.78, 3.0, 3.2, 0.32, size=16, bold=True, color=blue)
    add_text(slide, "Three layers of risk mitigation:", 0.78, 3.43, 4.5, 0.26, size=13.4, color=blue)
    risk_layers = [
        ("1", "Execution parity", "Live sequencing, broker-paper fills, and saved data must reproduce the post-session review path before capital is increased."),
        ("2", "Stress-based account sizing", "The sizing anchor is 3x historical intrabar MTM drawdown. A day that reaches the historical max stress level becomes a macro review day before continuation."),
        ("3", "Selective regime filters", "Known degradation zones, including unusually large opening ranges, are candidates for size reduction or exclusion if they preserve upside while reducing exposure."),
    ]
    y = 3.88
    for num, title, body in risk_layers:
        add_text(slide, f"{num})", 0.78, y, 0.35, 0.25, size=12.5, bold=True, color=dark)
        add_text(slide, title + ":", 1.12, y, 2.25, 0.25, size=12.5, bold=True, color=blue)
        add_text(slide, body, 3.15, y, 9.1, 0.42, size=11.3, color=dark)
        y += 0.76
    add_text(
        slide,
        "The objective of this stage is not to maximize advertised returns; it is to reduce execution risk to an acceptable level while preserving the historical edge profile.",
        0.78,
        6.42,
        11.45,
        0.28,
        size=12.1,
        bold=True,
        color=gold,
    )

    # 8. Two systems to test
    slide = prs.slides.add_slide(blank)
    add_header(slide, "Systems Selected For The Build/Test Runway", "gated and ungated intraday")
    benchmark_rows = funding_benchmark_rows(summaries)
    ppt_rows = [
        {
            "System": row["System"],
            "Norm Scale": row["Normalized Scale"],
            "Norm Net": row["Normalized Net"],
            "Norm Return": row["Normalized Return"],
            "Net/DD": row["Norm. Net/DD"],
            "$1M Books": row["$1M Books"],
            "$1M Net": row["$1M Net"],
            "$1M Return": row["$1M Return"],
        }
        for row in benchmark_rows
    ]
    add_table(
        slide,
        ppt_rows,
        ["System", "Norm Scale", "Norm Net", "Norm Return", "Net/DD", "$1M Books", "$1M Net", "$1M Return"],
        0.6,
        1.35,
        12.1,
        1.1,
        font_size=9.0,
        widths=[2.2, 1.25, 1.55, 1.25, 0.85, 1.1, 1.55, 1.15],
    )
    add_box(slide, 0.75, 3.1, 5.75, 2.15, radius=False)
    add_text(slide, "Flagship gated system", 1.08, 3.38, 4.75, 0.32, size=18, bold=True, color=blue)
    add_bullets(
        slide,
        [
            "Primary research candidate.",
            "More selective and historically more capital efficient.",
            "Requires stricter sequencing/tick proof before live confidence.",
        ],
        1.08,
        3.9,
        4.9,
        1.05,
        size=11.8,
    )
    add_box(slide, 6.9, 3.1, 5.75, 2.15, radius=False)
    add_text(slide, "Ungated price-action system", 7.22, 3.38, 4.75, 0.32, size=18, bold=True, color=blue)
    add_bullets(
        slide,
        [
            "Secondary system and operational benchmark.",
            "Simpler all-day rule family.",
            "Useful for testing cloud runtime, OCO lifecycle, fills, and reporting.",
        ],
        7.22,
        3.9,
        4.9,
        1.05,
        size=11.8,
    )
    add_text(slide, "Normalized rows use the same stress-capital anchor from the research tracker; fractional scale is comparison math, not live sizing.", 0.85, 6.08, 11.35, 0.25, size=10.8, color=slate)

    # 9. Founder partner alignment
    slide = prs.slides.add_slide(blank)
    add_header(slide, "Founder Partner Alignment", "why early support matters")
    add_text(
        slide,
        "Family and friends are taking the earliest calculated risk: funding an unproven live-test build before there is an audited managed-account record. That support creates the operating proof needed to build an auditable track record.",
        0.82,
        1.28,
        11.55,
        0.66,
        size=13.1,
        color=dark,
    )
    founder_rows = [
        {"Exchange": "Early risk capital", "Draft founder benefit": "Priority access to validated strategy capacity if the program launches."},
        {"Exchange": "Trust before audited performance", "Draft founder benefit": f"Target founder terms: {FOUNDER_PARTNER_MGMT_FEE:.0f}% management / {FOUNDER_PARTNER_PERFORMANCE_FEE:.0f}% performance fee."},
        {"Exchange": "Track-record runway", "Draft founder benefit": "Ability to carry a validated allocation forward or withdraw at defined review points."},
        {"Exchange": "Alignment signal", "Draft founder benefit": "Founder capital remains visibly aligned with the platform's long-term operating record."},
    ]
    add_table(slide, founder_rows, ["Exchange", "Draft founder benefit"], 0.78, 2.2, 11.75, 1.78, font_size=10.2, widths=[3.15, 8.6])
    add_box(slide, 0.95, 4.55, 5.45, 1.33, fill_color=RGBColor(246, 250, 253), radius=False)
    add_text(slide, "Founder account principle", 1.22, 4.78, 4.9, 0.28, size=15.2, bold=True, color=blue)
    add_text(
        slide,
        "Early supporters should not be treated like later outside allocators if the firm succeeds. The intended model is a documented founder class or side letter.",
        1.22,
        5.18,
        4.8,
        0.42,
        size=10.7,
        color=dark,
    )
    add_box(slide, 6.82, 4.55, 5.45, 1.33, fill_color=RGBColor(255, 249, 239), line_color=gold, radius=False)
    add_text(slide, "Compliance caveat", 7.1, 4.78, 4.9, 0.28, size=15.2, bold=True, color=gold)
    add_text(
        slide,
        "No fee, capacity, withdrawal, or account-management commitment is final until counsel-approved agreements exist.",
        7.1,
        5.2,
        4.8,
        0.42,
        size=10.8,
        color=dark,
    )

    # 10. Conservative performance framing
    slide = prs.slides.add_slide(blank)
    add_header(slide, "Conservative Performance Framing", "top 3 annual net years removed")
    perf_rows = []
    for row in summary_rows(summaries):
        perf_rows.append(
            {
                "Tier": row["Tier"],
                "System": row["System"].replace("Prior-Opposed Gated Intraday System", "Gated").replace("Ungated Intraday Breakout System", "Ungated"),
                "Market": row["Market"],
                "Min": row["Minimum"],
                "Window": row["Window"].replace(" to ", "-"),
                "Adv Avg/Yr": row["Advertised Avg/Yr"],
                "Adv Return/Yr": row["Advertised Return/Yr"],
                "Corr": row["Corr to QQQ"],
            }
        )
    add_table(
        slide,
        perf_rows,
        ["Tier", "System", "Market", "Min", "Window", "Adv Avg/Yr", "Adv Return/Yr", "Corr"],
        0.35,
        1.25,
        12.65,
        1.75,
        font_size=7.4,
        widths=[2.0, 1.2, 0.75, 1.0, 2.15, 1.3, 1.35, 0.65],
    )
    add_image_fit(slide, CHARTS / "average_annual_return_by_strategy.png", 0.92, 3.35, 11.4, 2.55)
    add_text(slide, "Full annual tables remain in the memo; the displayed average removes the three best net years before calculating the headline expected-return framing.", 0.95, 6.25, 11.1, 0.28, size=11, color=slate)

    # 11. Cost model
    slide = prs.slides.add_slide(blank)
    add_header(slide, f"{RUNWAY_MONTHS}-Month Funding And Cost Model", "StoneX schedule included")
    budget = pd.read_csv(EXHIBITS / BUDGET_EXHIBIT)
    budget_rows = [
        {"Category": r.Category, "Amount": money(float(r.Amount), 0)}
        for r in budget.itertuples(index=False)
        if r.Category
        in [
            "Founder development stipend",
            "Segregated MNQ research/test capital",
            "Broker/API/exchange/commission/slippage allowance",
            "NFA/Series 3/registration admin",
            "CTA disclosure/legal/compliance/audit reserve",
            "Contingency",
        ]
    ]
    add_table(slide, budget_rows, ["Category", "Amount"], 0.7, 1.22, 5.8, 2.15, font_size=9.5, widths=[4.25, 1.55])
    add_text(slide, f"Target raise: {money(STARTER_ASK, 0)}", 0.82, 3.72, 4.5, 0.36, size=22, bold=True, color=blue)
    add_text(slide, "The first runway funds engineering, data, CTA registration/disclosure prep, and MNQ test capital. NQ remains a future tier after execution evidence.", 0.82, 4.25, 5.4, 0.72, size=12.3, color=dark)
    stonex_rows = [
        {"Item": "Retail account minimum", "Fee": "$2,000"},
        {"Item": "Live simulation minimum", "Fee": "$5,000"},
        {"Item": "Standard futures", "Fee": "$1.29/side"},
        {"Item": "Micro futures", "Fee": "$0.50/side"},
        {"Item": "Simulation software", "Fee": "$49.95/month"},
    ]
    add_table(slide, stonex_rows, ["StoneX Item", "Fee"], 6.9, 1.22, 5.35, 2.15, font_size=10, widths=[3.35, 2.0])
    add_bullets(
        slide,
        [
            "Fee schedule is founder-supplied and requires account-opening confirmation.",
            "CTA registration/disclosure reserve is now a first-class cost before client trading.",
            "Commission-only model estimates are exported in the package exhibits.",
            "Exchange, NFA, clearing, market-data, margin-interest, wire, and slippage costs remain live-account assumptions.",
        ],
        6.98,
        4.05,
        5.0,
        1.42,
        size=11.9,
    )

    # 12. Twelve-month roadmap
    slide = prs.slides.add_slide(blank)
    add_header(slide, f"{RUNWAY_MONTHS}-Month Validation Roadmap", "evidence before scale")
    roadmap = [
        ("1-3", "Build The Rails", "CTA counsel kickoff, Series 3/NFA prep, source audit, CQG/StoneX adapter, feed shadowing, and replay parity."),
        ("4-6", "Paper The System", "Disclosure draft/review, MNQ broker-paper trial, restart drills, emergency flatten drills, slippage/cost audit, and first readiness review."),
        ("7-9", "Extend The Pilot", "Longer funded-paper or small-live MNQ observation, ungated comparison, investor reporting packet, and incident cadence."),
        ("10-12", "Decision Package", "Regime review, QQQ comparison refresh, compliance/accounting packet, registration/disclosure status, and final MNQ/NQ go/no-go."),
    ]
    x_positions = [0.65, 3.78, 6.91, 10.04]
    for (num, title, body), x in zip(roadmap, x_positions):
        add_box(slide, x, 1.42, 2.55, 3.45, fill_color=RGBColor(252, 253, 255), radius=False)
        add_text(slide, "Months", x + 0.18, 1.7, 0.72, 0.24, size=8.2, bold=True, color=RGBColor(122, 128, 137), align=PP_ALIGN.CENTER)
        add_text(slide, num, x + 0.95, 1.61, 0.8, 0.38, size=16.8, bold=True, color=blue, align=PP_ALIGN.CENTER)
        add_text(slide, title, x + 0.25, 2.35, 2.05, 0.44, size=11.2, bold=True, color=navy, align=PP_ALIGN.CENTER)
        add_text(slide, body, x + 0.23, 3.02, 2.08, 1.1, size=8.8, color=dark, align=PP_ALIGN.CENTER)
    add_text(
        slide,
        "Reporting templates cover equity/margin status, daily P&L, annual and monthly returns, open exposure, fill reconciliation, and operational incidents.",
        0.75,
        5.62,
        11.65,
        0.42,
        size=12.4,
        color=dark,
        align=PP_ALIGN.CENTER,
    )

    # 13. Founder bio
    slide = prs.slides.add_slide(blank)
    add_header(slide, "Founder Background", "Poe Kgengwenyane")
    add_box(slide, 0.75, 1.25, 4.15, 4.85, fill_color=RGBColor(252, 253, 255), radius=False)
    add_text(slide, "Market Research & Systems", 1.05, 1.58, 3.55, 0.3, size=16, bold=True, color=blue)
    add_bullets(
        slide,
        [
            "Four years of market research across forex, futures, and systematic trade design.",
            "Futures systems developer focused on broker-like replay, automated execution, and reporting.",
            "Current objective: convert research into a monitored, reviewable, multi-tenant platform.",
        ],
        1.05,
        2.1,
        3.45,
        1.75,
        size=11.2,
    )
    add_text(slide, "Contact", 1.05, 4.62, 1.3, 0.25, size=13.5, bold=True, color=blue)
    add_text(slide, "tshepok13@gmail.com\nlinkedin.com/in/poe-k-722858111", 1.05, 4.95, 3.4, 0.52, size=10.4, color=dark)
    add_box(slide, 5.3, 1.25, 7.2, 4.85, fill_color=white, radius=False)
    add_text(slide, "Relevant Operating Experience", 5.65, 1.58, 5.8, 0.3, size=16, bold=True, color=blue)
    add_bullets(
        slide,
        [
            "Security-critical distributed systems at Thales across TypeScript, Python, and C++.",
            "Automated test frameworks for REST APIs, event-driven microservices, OpenAPI, pytest, Jest/Supertest, and BDD/Behave.",
            "Backend platform engineering across Elixir/Phoenix, Ruby on Rails, Node/NestJS, AWS, OAuth, metrics, fraud controls, and service migrations.",
            "Quality automation and release testing experience in cybersecurity environments.",
            "Bachelor of Science in Computer Science; Bachelor's Degree in Biomedical and Mechanical Engineering.",
            "Practical infrastructure skills include Kubernetes, NoSQL, REST APIs, Docker, CI/CD, monitoring, and reporting.",
        ],
        5.65,
        2.1,
        6.3,
        3.38,
        size=10.6,
    )

    prs.save(POWERPOINT_PATH)


def readme(summaries: Dict[str, Dict[str, object]]) -> str:
    petit = summaries["mnq_prior_opposed"]
    grand = summaries["nq_prior_opposed_full"]
    return "\n".join(
        [
            "# Family/Friends Funding Package",
            "",
            f"Private Markdown package for the **{PROGRAM_NAME}** build/test runway.",
            "",
            "## Documents",
            "",
            "- [PITCH_DECK.md](PITCH_DECK.md)",
            "- [SHORT_MEMO.md](SHORT_MEMO.md)",
            "- [COST_AND_RUNWAY.md](COST_AND_RUNWAY.md)",
            "- [REPORTING_TEMPLATES.md](REPORTING_TEMPLATES.md)",
            "- [RISK_AND_DISCLOSURE_NOTES.md](RISK_AND_DISCLOSURE_NOTES.md)",
            "- [VALIDATION.md](VALIDATION.md)",
            "",
            "## PowerPoint Draft",
            "",
            "- [PITCH_DECK.pptx](PITCH_DECK.pptx)",
            "",
            "## PDF Drafts",
            "",
            "- [pdfs/PITCH_DECK.pdf](pdfs/PITCH_DECK.pdf)",
            "- [pdfs/SHORT_MEMO.pdf](pdfs/SHORT_MEMO.pdf)",
            "- [pdfs/COST_AND_RUNWAY.pdf](pdfs/COST_AND_RUNWAY.pdf)",
            "- [pdfs/REPORTING_TEMPLATES.pdf](pdfs/REPORTING_TEMPLATES.pdf)",
            "- [pdfs/RISK_AND_DISCLOSURE_NOTES.pdf](pdfs/RISK_AND_DISCLOSURE_NOTES.pdf)",
            "",
            "## Key Read",
            "",
            f"- Ask: **about {money(STARTER_ASK, 0)}** for a CTA-registration-first {RUNWAY_MONTHS}-month build/test runway.",
            f"- First live-test path: **{PETIT_NAME}**, {money(MNQ_MINIMUM, 0)} minimum, MNQ-focused.",
            f"- Future larger tier: **{GRAND_NAME}**, {money(NQ_MINIMUM, 0)} minimum, NQ-focused.",
            f"- {PETIT_NAME} conservative advertised average after removing top 3 annual net years: **{money(float(petit['advertised_avg_net']), 0)} / {pct(float(petit['advertised_avg_return_pct']), 1)} per year**.",
            f"- {GRAND_NAME} conservative advertised average after removing top 3 annual net years: **{money(float(grand['advertised_avg_net']), 0)} / {pct(float(grand['advertised_avg_return_pct']), 1)} per year**.",
            f"- Draft founder alignment principle: early family/friends accounts target **{FOUNDER_PARTNER_MGMT_FEE:.0f}% management / {FOUNDER_PARTNER_PERFORMANCE_FEE:.0f}% performance fee** if a compliant future offering is launched, subject to counsel-approved documents.",
            "- QQQ benchmarks invest the full tier minimum over the same clipped strategy window, matching the fair benchmark convention from the research scripts.",
            "- All performance is hypothetical/backtested and unaudited.",
            "",
            "## Exhibits",
            "",
            "- `exhibits/strategy_summary.csv`",
            "- `exhibits/annual_returns_vs_qqq.csv`",
            f"- `exhibits/{BUDGET_EXHIBIT}`",
            "- `exhibits/regulatory_cost_assumptions.csv`",
            "- `exhibits/stonex_fee_estimate.csv`",
            "- `exhibits/scaling_10y_nq_prior_vs_qqq_50k_summary.csv`",
            "- `exhibits/scaling_10y_nq_prior_vs_qqq_50k_daily.csv`",
            "- `exhibits/source_manifest_internal.csv`",
            "",
            "## Charts",
            "",
            "- `charts/scaling_10y_nq_prior_vs_qqq_50k.png`",
            "- `charts/preferred_tiers_vs_qqq_equity.png`",
            "- `charts/nq_prior_opposed_full_vs_qqq_equity.png`",
            "- `charts/mini_annual_returns_vs_qqq.png`",
            "- `charts/average_annual_return_by_strategy.png`",
            "- `charts/daily_correlation_to_qqq.png`",
        ]
    )


def markdown_to_pdf(md_path: Path, pdf_path: Path) -> None:
    """Small dependency-free Markdown-to-PDF renderer for draft review packets."""

    width, height = 8.5, 11.0
    left, right = 0.07, 0.94
    top, bottom = 0.955, 0.055
    line_h = 0.022

    def new_page(pdf: PdfPages):
        fig = plt.figure(figsize=(width, height))
        fig.patch.set_facecolor("white")
        ax = fig.add_axes([0, 0, 1, 1])
        ax.axis("off")
        return fig, ax, top

    def finish_page(pdf: PdfPages, fig):
        pdf.savefig(fig, dpi=160)
        plt.close(fig)

    def draw_wrapped(ax, x, y, text, size=9, weight="normal", family="DejaVu Sans", indent=0.0):
        clean = re.sub(r"\*\*(.*?)\*\*", r"\1", text).replace("`", "")
        chars = max(35, int(98 - indent * 150 - max(size - 9, 0) * 7))
        wrapped = textwrap.wrap(clean, width=chars) or [""]
        for idx, part in enumerate(wrapped):
            ax.text(
                x + indent,
                y,
                part,
                ha="left",
                va="top",
                fontsize=size,
                fontweight=weight,
                family=family,
                color="#111111",
            )
            y -= line_h * (size / 9.0)
        return y

    def ensure_space(pdf: PdfPages, fig, ax, y, needed):
        if y - needed >= bottom:
            return fig, ax, y
        finish_page(pdf, fig)
        return new_page(pdf)

    image_pattern = re.compile(r"!\[[^\]]*\]\(([^)]+)\)")

    with PdfPages(pdf_path) as pdf:
        fig, ax, y = new_page(pdf)
        for raw_line in md_path.read_text(encoding="utf-8").splitlines():
            line = raw_line.rstrip()
            image_match = image_pattern.match(line.strip())
            if image_match:
                image_path = (md_path.parent / image_match.group(1)).resolve()
                if image_path.exists():
                    fig, ax, y = ensure_space(pdf, fig, ax, y, 0.36)
                    img = plt.imread(image_path)
                    aspect = img.shape[1] / img.shape[0]
                    box_w = right - left
                    box_h = min(0.33, box_w / aspect * (width / height))
                    image_ax = fig.add_axes([left, y - box_h, box_w, box_h])
                    image_ax.imshow(img)
                    image_ax.axis("off")
                    y -= box_h + 0.025
                continue

            if not line:
                y -= line_h * 0.55
                if y < bottom:
                    finish_page(pdf, fig)
                    fig, ax, y = new_page(pdf)
                continue

            if line.strip() == "---":
                fig, ax, y = ensure_space(pdf, fig, ax, y, 0.03)
                ax.plot([left, right], [y, y], color="#c9c9c9", linewidth=0.8)
                y -= line_h
                continue

            if line.startswith("# "):
                fig, ax, y = ensure_space(pdf, fig, ax, y, 0.06)
                y = draw_wrapped(ax, left, y, line[2:], size=15, weight="bold")
                y -= line_h * 0.35
                continue
            if line.startswith("## "):
                fig, ax, y = ensure_space(pdf, fig, ax, y, 0.045)
                y = draw_wrapped(ax, left, y, line[3:], size=12, weight="bold")
                y -= line_h * 0.18
                continue

            if line.startswith("|"):
                fig, ax, y = ensure_space(pdf, fig, ax, y, 0.028)
                if re.fullmatch(r"\|\s*:?-+:?\s*(\|\s*:?-+:?\s*)+\|?", line):
                    continue
                y = draw_wrapped(ax, left, y, line, size=6.8, family="DejaVu Sans Mono")
                continue

            if line.startswith("- "):
                fig, ax, y = ensure_space(pdf, fig, ax, y, 0.035)
                y = draw_wrapped(ax, left, y, "• " + line[2:], size=9, indent=0.015)
                continue

            if re.match(r"^\d+\. ", line):
                fig, ax, y = ensure_space(pdf, fig, ax, y, 0.035)
                y = draw_wrapped(ax, left, y, line, size=9, indent=0.015)
                continue

            fig, ax, y = ensure_space(pdf, fig, ax, y, 0.035)
            y = draw_wrapped(ax, left, y, line, size=9)

        finish_page(pdf, fig)


def write_pdfs(doc_names: Optional[List[str]] = None) -> None:
    PDF_DIR.mkdir(parents=True, exist_ok=True)
    for doc_name in (doc_names or ["README.md"] + PUBLIC_DOCS + ["VALIDATION.md"]):
        markdown_to_pdf(OUT / doc_name, PDF_DIR / f"{Path(doc_name).stem}.pdf")


def validate_public_docs() -> str:
    missing_links = []
    link_pattern = re.compile(r"\[[^\]]+\]\(([^)]+)\)")
    for doc_name in ["README.md"] + PUBLIC_DOCS:
        path = OUT / doc_name
        text = path.read_text(encoding="utf-8")
        for link in link_pattern.findall(text):
            if link.startswith(("http://", "https://", "mailto:")):
                continue
            target = (path.parent / link).resolve()
            if "#" in link:
                target = (path.parent / link.split("#", 1)[0]).resolve()
            if not target.exists():
                missing_links.append(f"{doc_name}: {link}")

    leakage_terms = [
        "ST+PMC",
        "S_1_1_3",
        "v2b_prior_opposed",
        "entry_ts",
        "trigger level",
        "hourly gate",
    ]
    leakage_hits = []
    for doc_name in PUBLIC_DOCS:
        text = (OUT / doc_name).read_text(encoding="utf-8")
        for term in leakage_terms:
            if term.lower() in text.lower():
                leakage_hits.append(f"{doc_name}: {term}")

    strategy_summary = pd.read_csv(EXHIBITS / "strategy_summary.csv")
    avg_checks = []
    for row in strategy_summary.itertuples(index=False):
        src = next(s for s in SOURCES if s.slug == row.slug)
        yearly = read_yearly(src)
        expected_net, expected_pct, _, _ = conservative_average(yearly)
        ok = math.isclose(float(row.advertised_avg_net), expected_net, rel_tol=0, abs_tol=0.01) and math.isclose(
            float(row.advertised_avg_return_pct), expected_pct, rel_tol=0, abs_tol=0.0001
        )
        avg_checks.append(
            {
                "Strategy": row.slug,
                "Top3Removed": row.removed_top3_years,
                "AdvertisedAvgNet": money(float(row.advertised_avg_net), 0),
                "AdvertisedAvgReturn": pct(float(row.advertised_avg_return_pct), 1),
                "Pass": "yes" if ok else "no",
            }
        )

    return "\n".join(
        [
            "# Validation",
            "",
            "Generated by `scripts/generate_family_friends_funding_package.py`.",
            "",
            "## Link Check",
            "",
            f"- Missing local links: **{len(missing_links)}**",
            "",
            "\n".join(f"- {m}" for m in missing_links) if missing_links else "No missing local links found.",
            "",
            "## Public Disclosure Leakage Scan",
            "",
            f"- Hits: **{len(leakage_hits)}**",
            "",
            "\n".join(f"- {h}" for h in leakage_hits) if leakage_hits else "No configured leakage terms found in public docs.",
            "",
            "## Advertised Average Check",
            "",
            md_table(avg_checks, ["Strategy", "Top3Removed", "AdvertisedAvgNet", "AdvertisedAvgReturn", "Pass"]),
            "",
            "## Data Freshness Note",
            "",
            "QQQ benchmark equity is reset to the full tier minimum on the first matched QQQ trading day inside each strategy's actual replay window, matching the fair-benchmark convention used by `scripts/fair_benchmark_comparison.py`. The NQ full-history raw run validates through early March 2026 because the daily regime and gate support files end there, even though the restored raw DBN extends later.",
        ]
    )


def main() -> None:
    ensure_dirs()
    summaries = {src.slug: summarize_strategy(src) for src in SOURCES}
    write_exhibits(summaries)
    write_charts(summaries)

    (OUT / "README.md").write_text(readme(summaries), encoding="utf-8")
    (OUT / "PITCH_DECK.md").write_text(pitch_deck(summaries), encoding="utf-8")
    (OUT / "SHORT_MEMO.md").write_text(short_memo(summaries), encoding="utf-8")
    (OUT / "COST_AND_RUNWAY.md").write_text(cost_and_runway(), encoding="utf-8")
    (OUT / "REPORTING_TEMPLATES.md").write_text(reporting_templates(), encoding="utf-8")
    (OUT / "RISK_AND_DISCLOSURE_NOTES.md").write_text(risk_and_disclosure_notes(), encoding="utf-8")
    (OUT / "VALIDATION.md").write_text("", encoding="utf-8")
    write_pdfs(["README.md"] + PUBLIC_DOCS)
    write_powerpoint(summaries)
    (OUT / "VALIDATION.md").write_text(validate_public_docs(), encoding="utf-8")
    write_pdfs(["VALIDATION.md"])
    write_run_manifest(
        OUT,
        output_paths=[
            OUT / "README.md",
            OUT / "PITCH_DECK.md",
            OUT / "SHORT_MEMO.md",
            OUT / "COST_AND_RUNWAY.md",
            OUT / "REPORTING_TEMPLATES.md",
            OUT / "RISK_AND_DISCLOSURE_NOTES.md",
            OUT / "VALIDATION.md",
        ],
        strategy_config={"driver": "generate_family_friends_funding_package"},
        causality_mode="audit",
        extra={"source_slugs": [src.slug for src in SOURCES]},
        repo_root=ROOT,
    )
    print(f"Wrote {OUT}")


if __name__ == "__main__":
    main()
